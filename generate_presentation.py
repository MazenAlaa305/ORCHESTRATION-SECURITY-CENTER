#!/usr/bin/env python3
"""
Orchestration Security Center — Graduation Project Defense Presentation Generator
==============================================================
Generates a 30-slide academic defense PPTX with a "Futuristic Glass / iOS" aesthetic.

Theme palette
─────────────
  BG_DARK       #0A0E27   (deep space blue)
  BG_CARD       #141A36   (glass card fill)
  BG_CARD_ALT   #1A2240   (lighter glass variant)
  ACCENT_CYAN   #00D4FF   (primary accent)
  ACCENT_BLUE   #4A90D9   (secondary accent)
  ACCENT_GREEN  #00FF88   (safe / success)
  ACCENT_ORANGE #FFAA00   (warning)
  ACCENT_RED    #FF0055   (critical / danger)
  TEXT_WHITE    #FFFFFF
  TEXT_LIGHT    #C0C8D8   (body copy)
  TEXT_DIM      #7B8BA8   (captions)

Typography: Merriweather Light (all elements)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
import math
import os

# ── Colour Palette ──────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x0A, 0x0E, 0x27)
BG_CARD       = RGBColor(0x14, 0x1A, 0x36)
BG_CARD_ALT   = RGBColor(0x1A, 0x22, 0x40)
ACCENT_CYAN   = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_BLUE   = RGBColor(0x4A, 0x90, 0xD9)
ACCENT_GREEN  = RGBColor(0x00, 0xFF, 0x88)
ACCENT_ORANGE = RGBColor(0xFF, 0xAA, 0x00)
ACCENT_RED    = RGBColor(0xFF, 0x00, 0x55)
TEXT_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT    = RGBColor(0xC0, 0xC8, 0xD8)
TEXT_DIM      = RGBColor(0x7B, 0x8B, 0xA8)
BORDER_GLOW   = RGBColor(0x00, 0xD4, 0xFF)  # cyan glow lines

FONT_NAME = "Merriweather Light"
FONT_FALLBACK = "Merriweather"

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ─────────────────────────────────────────────────────────────────

def _hex(r, g, b):
    return RGBColor(r, g, b)

def _set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=None, border_width=Pt(1), radius=None):
    """Add a rounded rectangle 'glass card'."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Transparency for glass effect (20%)
    shape.fill.fore_color.brightness = 0.0
    return shape

def _tf(shape):
    """Return the text frame of a shape, ready for writing."""
    tf = shape.text_frame
    tf.word_wrap = True
    return tf

def _set_font(run, size=Pt(14), color=TEXT_WHITE, bold=False, italic=False, name=FONT_NAME):
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name

def _add_text_box(slide, left, top, width, height, text, size=Pt(14), color=TEXT_WHITE,
                  bold=False, alignment=PP_ALIGN.LEFT, line_space=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, color=color, bold=bold)
    p.space_after = Pt(2)
    return txBox

def _add_title(slide, text, subtitle=None, left=Inches(0.8), top=Inches(0.35)):
    """Add a slide title with optional subtitle and accent line."""
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.0), Inches(0.6), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_CYAN
    line.line.fill.background()

    _add_text_box(slide, left + Inches(0.75), top - Inches(0.15), Inches(10), Inches(0.6),
                  text, size=Pt(28), color=TEXT_WHITE, bold=True)
    if subtitle:
        _add_text_box(slide, left + Inches(0.75), top + Inches(0.45), Inches(10), Inches(0.4),
                      subtitle, size=Pt(14), color=TEXT_DIM)

def _add_bullet_card(slide, left, top, width, height, title, bullets, title_color=ACCENT_CYAN,
                     bullet_color=TEXT_LIGHT, border=ACCENT_BLUE):
    card = _add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=border, border_width=Pt(1))
    tf = _tf(card)
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.15)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_font(run, size=Pt(16), color=title_color, bold=True)
    p.space_after = Pt(8)

    for b in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = f"  {b}"
        _set_font(run, size=Pt(12), color=bullet_color)
    return card

def _add_kpi_box(slide, left, top, width, height, label, value, color=ACCENT_CYAN):
    card = _add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=color, border_width=Pt(1.5))
    tf = _tf(card)
    tf.margin_top = Inches(0.15)
    tf.margin_left = Inches(0.1)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(value)
    _set_font(run, size=Pt(32), color=color, bold=True)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    _set_font(run2, size=Pt(11), color=TEXT_DIM)
    return card

def _add_flow_arrow(slide, left, top, width=Inches(0.5), color=ACCENT_CYAN):
    """Draw a right-pointing arrow."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow

def _add_down_arrow(slide, left, top, height=Inches(0.4), color=ACCENT_CYAN):
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.35), height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow

def _slide_number(slide, num):
    """Add slide number in bottom-right."""
    _add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1), Inches(0.4),
                  str(num), size=Pt(10), color=TEXT_DIM, alignment=PP_ALIGN.RIGHT)

def _add_footer_bar(slide):
    """Thin cyan accent bar at page bottom."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.3), SLIDE_W, Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_CYAN
    bar.line.fill.background()

def new_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide)
    _add_footer_bar(slide)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
# Large centered title
_add_text_box(s, Inches(1), Inches(1.5), Inches(11.3), Inches(1),
              "ORCHESTRATION SECURITY CENTER", size=Pt(60), color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
_add_text_box(s, Inches(1), Inches(2.6), Inches(11.3), Inches(0.8),
              "Autonomous Security Orchestration Platform for SMEs",
              size=Pt(24), color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

# Accent line
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.5), Inches(2.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_CYAN
line.line.fill.background()

# Meta info
_add_text_box(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
              "Graduation Project Defense  |  2026",
              size=Pt(16), color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

_add_text_box(s, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
              "AI-Driven DAST  \u2022  Multi-Agent Pipeline  \u2022  Unified Risk Engine  \u2022  SIEM/SOAR Integration",
              size=Pt(13), color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Decorative corner shapes
for (cx, cy) in [(Inches(0.3), Inches(0.3)), (Inches(12.5), Inches(0.3)),
                  (Inches(0.3), Inches(6.8)), (Inches(12.5), Inches(6.8))]:
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.5), Inches(0.5))
    d.fill.solid()
    d.fill.fore_color.rgb = ACCENT_BLUE
    d.line.fill.background()
    # make it semi-transparent via brightness hack
    d.fill.fore_color.brightness = 0.6

_slide_number(s, 1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "AGENDA", "Presentation Roadmap")

agenda_items = [
    ("01", "Problem Statement & Motivation"),
    ("02", "Research Objectives"),
    ("03", "Literature Review & Gaps"),
    ("04", "System Architecture Overview"),
    ("05", "Microservices Decomposition"),
    ("06", "Multi-Agent Orchestration Pipeline"),
    ("07", "Unified Risk Engine — Deterministic Scoring"),
    ("08", "Deterministic vs. Non-Deterministic Logic"),
    ("09", "SIEM / SOAR Integration"),
    ("10", "Living Lab Environment"),
    ("11", "Dashboard & Real-Time Visualization"),
    ("12", "Results, Evaluation & Novel Contributions"),
    ("13", "Future Work & Conclusion"),
]

col1 = agenda_items[:7]
col2 = agenda_items[7:]

for i, (num, item) in enumerate(col1):
    y = Inches(1.3) + Inches(i * 0.7)
    _add_text_box(s, Inches(1.0), y, Inches(0.6), Inches(0.5), num, size=Pt(20), color=ACCENT_CYAN, bold=True)
    _add_text_box(s, Inches(1.7), y + Inches(0.03), Inches(5), Inches(0.5), item, size=Pt(15), color=TEXT_LIGHT)

for i, (num, item) in enumerate(col2):
    y = Inches(1.3) + Inches(i * 0.7)
    _add_text_box(s, Inches(7.2), y, Inches(0.6), Inches(0.5), num, size=Pt(20), color=ACCENT_CYAN, bold=True)
    _add_text_box(s, Inches(7.9), y + Inches(0.03), Inches(5), Inches(0.5), item, size=Pt(15), color=TEXT_LIGHT)

_slide_number(s, 2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ABSTRACT
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "ABSTRACT", "Executive Summary")

abstract_text = (
    "Orchestration Security Center is an AI-driven Dynamic Application Security Testing (DAST) platform "
    "purpose-built for Small and Medium Enterprises. The system orchestrates a multi-agent "
    "autonomous pipeline\u2014ReconAgent, AttackAgent, ValidationAgent, and ReportingAgent\u2014"
    "powered by Google Gemini LLM, unifying industry-standard tools (Nmap, Nuclei, OpenVAS, "
    "Wazuh) under a single microservices architecture.\n\n"
    "The platform introduces a Unified Risk Engine that computes deterministic risk scores "
    "by weighting vulnerability severity, asset criticality, and network exposure. A living "
    "lab with 4 subnets and 11 vulnerable containers simulates realistic SME environments "
    "for validation. Real-time dashboards, SIEM/SOAR integration, and non-technical advisory "
    "generation bridge the gap between complex security telemetry and actionable business decisions."
)

card = _add_shape(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(4.5),
                  fill_color=BG_CARD, border_color=ACCENT_BLUE, border_width=Pt(1))
tf = _tf(card)
tf.margin_left = Inches(0.4)
tf.margin_top = Inches(0.3)
tf.margin_right = Inches(0.4)
p = tf.paragraphs[0]
run = p.add_run()
run.text = abstract_text
_set_font(run, size=Pt(15), color=TEXT_LIGHT)
p.line_spacing = Pt(24)

# Keywords
_add_text_box(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
              "Keywords:  DAST  \u2022  Multi-Agent Systems  \u2022  Risk Quantification  \u2022  SIEM/SOAR  \u2022  LLM-Enhanced Security  \u2022  SME Cybersecurity",
              size=Pt(12), color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

_slide_number(s, 3)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "PROBLEM STATEMENT", "The SME Security Gap")

problems = [
    ("\u26A0  Tool Fragmentation",
     "SMEs lack resources to manage Nmap, Nuclei, OpenVAS, and Wazuh independently. "
     "Each tool has its own interface, output format, and expertise requirements."),
    ("\u26A0  Alert Fatigue",
     "Security tools generate thousands of findings with no unified prioritization. "
     "Non-technical stakeholders cannot interpret raw CVE data or CVSS scores."),
    ("\u26A0  Reactive Posture",
     "Most SMEs discover vulnerabilities after exploitation. Continuous automated testing "
     "and real-time risk dashboards are typically enterprise-only capabilities."),
    ("\u26A0  Expertise Shortage",
     "There is a global deficit of 3.4M cybersecurity professionals (ISC\u00B2 2023). "
     "SMEs cannot hire dedicated security teams or afford managed SIEM services."),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(1.2) + Inches(i * 1.45)
    card = _add_shape(s, Inches(0.8), y, Inches(11.7), Inches(1.3),
                      fill_color=BG_CARD, border_color=ACCENT_RED if i < 2 else ACCENT_ORANGE, border_width=Pt(1))
    tf = _tf(card)
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_font(run, size=Pt(16), color=ACCENT_RED if i < 2 else ACCENT_ORANGE, bold=True)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = desc
    _set_font(run2, size=Pt(12), color=TEXT_LIGHT)
    p2.space_before = Pt(6)

_slide_number(s, 4)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RESEARCH OBJECTIVES
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "RESEARCH OBJECTIVES", "Goals & Deliverables")

objectives = [
    ("RO-1", "Design a multi-agent autonomous DAST pipeline that chains reconnaissance, attack, validation, and reporting stages without human intervention."),
    ("RO-2", "Develop a Unified Risk Engine with deterministic scoring that weights vulnerability severity against asset criticality and network exposure."),
    ("RO-3", "Integrate heterogeneous security tools (Nmap, Nuclei, OpenVAS, Wazuh) via a unified microservices API layer."),
    ("RO-4", "Implement AI-enhanced advisory generation using Google Gemini LLM to translate technical findings into non-technical business language."),
    ("RO-5", "Build a living lab simulation environment modeling a realistic 4-subnet SME network topology for validation."),
    ("RO-6", "Deliver a real-time dashboard with WebSocket-driven updates and SIEM/SOAR integration for operational visibility."),
]

for i, (code, desc) in enumerate(objectives):
    y = Inches(1.2) + Inches(i * 0.95)
    # Code badge
    badge = _add_shape(s, Inches(0.8), y, Inches(1.0), Inches(0.5),
                       fill_color=ACCENT_CYAN, border_color=None)
    tf = _tf(badge)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = code
    _set_font(run, size=Pt(13), color=BG_DARK, bold=True)
    # Description
    _add_text_box(s, Inches(2.0), y, Inches(10.5), Inches(0.8),
                  desc, size=Pt(13), color=TEXT_LIGHT)

_slide_number(s, 5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LITERATURE REVIEW
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "LITERATURE REVIEW & RESEARCH GAPS", "Related Work Analysis")

# Comparison table as cards
tools_data = [
    ("OWASP ZAP", "Open-source DAST", "Single-tool, no orchestration, no AI advisory", ACCENT_ORANGE),
    ("Burp Suite Pro", "Commercial DAST", "Expensive, manual-heavy, no SIEM integration", ACCENT_ORANGE),
    ("DefectDojo", "Vuln aggregation", "No active scanning, no multi-agent pipeline", ACCENT_ORANGE),
    ("Qualys VMDR", "Enterprise VM", "Cost-prohibitive for SMEs, cloud-only", ACCENT_RED),
    ("Orchestration Security Center", "AI-Driven DAST + SOAR", "Unified pipeline, deterministic scoring, LLM advisory", ACCENT_GREEN),
]

# Column headers
for j, h in enumerate(["Platform", "Category", "Limitation / Advantage"]):
    x = Inches(0.8) + Inches(j * 4.0)
    _add_text_box(s, x, Inches(1.15), Inches(3.8), Inches(0.4),
                  h, size=Pt(13), color=ACCENT_CYAN, bold=True)

for i, (name, cat, note, clr) in enumerate(tools_data):
    y = Inches(1.65) + Inches(i * 0.95)
    row_bg = BG_CARD if i % 2 == 0 else BG_CARD_ALT
    card = _add_shape(s, Inches(0.8), y, Inches(11.7), Inches(0.8),
                      fill_color=row_bg, border_color=clr if i == 4 else None, border_width=Pt(1.5) if i == 4 else Pt(0))
    # Name
    _add_text_box(s, Inches(1.0), y + Inches(0.15), Inches(3.5), Inches(0.5),
                  name, size=Pt(14), color=TEXT_WHITE, bold=(i == 4))
    # Category
    _add_text_box(s, Inches(4.8), y + Inches(0.15), Inches(3.5), Inches(0.5),
                  cat, size=Pt(13), color=TEXT_LIGHT)
    # Note
    _add_text_box(s, Inches(8.8), y + Inches(0.15), Inches(3.5), Inches(0.5),
                  note, size=Pt(12), color=clr)

_slide_number(s, 6)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SYSTEM ARCHITECTURE OVERVIEW (High-Level)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "SYSTEM ARCHITECTURE OVERVIEW", "Microservices & Integration Points")

# Three-tier visualization
tiers = [
    ("PRESENTATION TIER", Inches(0.8), ACCENT_CYAN,
     ["React/Vite SPA (Port 5173)", "WebSocket Client (/ws/logs)", "Real-Time Dashboard", "RiskScore Gauge, Agent Logs"]),
    ("APPLICATION TIER", Inches(3.3), ACCENT_BLUE,
     ["FastAPI Gateway (Port 8000)", "Agent Orchestrator (4 Agents)", "Celery Workers + Beat", "Unified Risk Engine"]),
    ("DATA & SECURITY TIER", Inches(5.8), ACCENT_GREEN,
     ["PostgreSQL (Port 5432)", "Redis Pub/Sub (Port 6379)", "Elasticsearch/Kibana SIEM", "Wazuh Manager + n8n SOAR"]),
]

for label, y, color, items in tiers:
    card = _add_shape(s, Inches(0.8), y, Inches(11.7), Inches(2.2),
                      fill_color=BG_CARD, border_color=color, border_width=Pt(1.5))
    # Tier label
    badge = _add_shape(s, Inches(1.0), y + Inches(0.15), Inches(2.8), Inches(0.45),
                       fill_color=color)
    tf = _tf(badge)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    _set_font(run, size=Pt(13), color=BG_DARK, bold=True)

    # Items in row
    for j, item in enumerate(items):
        x = Inches(1.0) + Inches(j * 2.9)
        item_card = _add_shape(s, x, y + Inches(0.8), Inches(2.7), Inches(1.1),
                               fill_color=BG_CARD_ALT, border_color=color, border_width=Pt(0.5))
        tf = _tf(item_card)
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = item
        _set_font(run, size=Pt(11), color=TEXT_LIGHT)

# Arrows between tiers
_add_down_arrow(s, Inches(6.4), Inches(3.05), Inches(0.25))
_add_down_arrow(s, Inches(6.4), Inches(5.55), Inches(0.25))

_slide_number(s, 7)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MICROSERVICES DECOMPOSITION
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "MICROSERVICES DECOMPOSITION", "11 Containerized Services")

services = [
    ("Backend (FastAPI)", "8000", "API + WebSocket gateway", ACCENT_CYAN),
    ("Frontend (Vite)", "5173", "React SPA dashboard", ACCENT_CYAN),
    ("PostgreSQL", "5432", "Primary data store", ACCENT_BLUE),
    ("Redis", "6379", "Celery broker + Pub/Sub", ACCENT_BLUE),
    ("Celery Worker", "\u2014", "Async task processing", ACCENT_BLUE),
    ("Celery Beat", "\u2014", "Hourly scan scheduler", ACCENT_BLUE),
    ("OpenVAS/GVM", "9390", "Vulnerability scanner", ACCENT_GREEN),
    ("Elasticsearch", "9200", "SIEM log aggregation", ACCENT_GREEN),
    ("Kibana", "5601", "SIEM visualization", ACCENT_GREEN),
    ("Wazuh Manager", "55000", "Security event mgmt", ACCENT_ORANGE),
    ("n8n", "5678", "SOAR orchestration", ACCENT_ORANGE),
]

# Header
for j, h in enumerate(["Service", "Port", "Role"]):
    x = [Inches(0.8), Inches(4.5), Inches(6.0)][j]
    w = [Inches(3.5), Inches(1.3), Inches(6.0)][j]
    _add_text_box(s, x, Inches(1.1), w, Inches(0.35),
                  h, size=Pt(12), color=ACCENT_CYAN, bold=True)

for i, (name, port, role, clr) in enumerate(services):
    y = Inches(1.55) + Inches(i * 0.5)
    row_bg = BG_CARD if i % 2 == 0 else BG_CARD_ALT
    _add_shape(s, Inches(0.8), y, Inches(11.7), Inches(0.45), fill_color=row_bg)
    # Colored dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.1), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = clr
    dot.line.fill.background()

    _add_text_box(s, Inches(1.2), y + Inches(0.05), Inches(3.0), Inches(0.35),
                  name, size=Pt(12), color=TEXT_WHITE)
    _add_text_box(s, Inches(4.5), y + Inches(0.05), Inches(1.3), Inches(0.35),
                  port, size=Pt(12), color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    _add_text_box(s, Inches(6.0), y + Inches(0.05), Inches(6.0), Inches(0.35),
                  role, size=Pt(12), color=TEXT_LIGHT)

_slide_number(s, 8)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DATABASE SCHEMA
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "DATABASE ARCHITECTURE", "SQLAlchemy ORM — 9 Core Tables")

tables = [
    ("Target", "UUID PK, name, base_url, source, tech_stack (JSON),\nasset_value, data_sensitivity, auth_method"),
    ("Scan", "UUID PK, target_id FK, status, scan_type,\nrisk_score (0-100), health_score, agent_thoughts (JSON)"),
    ("Vulnerability", "UUID PK, scan_id FK, type, severity, status,\nconfidence_score, ai_validation_result (JSON), PoC, remediation"),
    ("AgentLog", "UUID PK, scan_id FK, agent_name, action,\nreasoning (JSON), input_data, output_data"),
    ("Endpoint", "UUID PK, target_id FK, url, method,\nparameters (JSON), auth_required"),
    ("ScanAsset", "INT PK, scan_id FK, ip_address, hostname,\nos_name, device_type, ai_insight (JSON)"),
    ("AssetService", "INT PK, asset_id FK, port, protocol,\nservice_name, product, version, CPE"),
    ("NetworkAsset", "INT PK, ip_address UNIQUE, status,\ncriticality, risk_score, first/last_seen"),
    ("ActionItem", "INT PK, scan_id FK, title, priority,\nstatus, type (REMEDIATION | CONFIGURATION)"),
]

col1 = tables[:5]
col2 = tables[5:]

for i, (name, fields) in enumerate(col1):
    y = Inches(1.2) + Inches(i * 1.2)
    card = _add_shape(s, Inches(0.6), y, Inches(6.0), Inches(1.05),
                      fill_color=BG_CARD, border_color=ACCENT_BLUE, border_width=Pt(1))
    tf = _tf(card)
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name
    _set_font(run, size=Pt(13), color=ACCENT_CYAN, bold=True)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = fields
    _set_font(run2, size=Pt(10), color=TEXT_DIM)

for i, (name, fields) in enumerate(col2):
    y = Inches(1.2) + Inches(i * 1.2)
    card = _add_shape(s, Inches(6.9), y, Inches(6.0), Inches(1.05),
                      fill_color=BG_CARD, border_color=ACCENT_BLUE, border_width=Pt(1))
    tf = _tf(card)
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name
    _set_font(run, size=Pt(13), color=ACCENT_CYAN, bold=True)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = fields
    _set_font(run2, size=Pt(10), color=TEXT_DIM)

_slide_number(s, 9)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — MULTI-AGENT PIPELINE OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "MULTI-AGENT ORCHESTRATION PIPELINE", "4-Stage Autonomous DAST Workflow")

agents = [
    ("1", "RECON\nAGENT", "Nmap scan\nPlaywright crawl\nTech stack detect\nEndpoint discovery", ACCENT_BLUE),
    ("2", "ATTACK\nAGENT", "Nuclei templates\nPayload injection\nService heuristics\nConcurrent testing", ACCENT_RED),
    ("3", "VALIDATION\nAGENT", "Confidence filter >0.6\nLLM re-validation\nFalse positive detect\nDeterministic rules", ACCENT_ORANGE),
    ("4", "REPORTING\nAGENT", "Executive summary\nPoC generation\nRemediation templates\nMarkdown report", ACCENT_GREEN),
]

for i, (num, name, desc, color) in enumerate(agents):
    x = Inches(0.5) + Inches(i * 3.2)
    # Agent card
    card = _add_shape(s, x, Inches(1.4), Inches(2.7), Inches(4.5),
                      fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.0), Inches(1.6), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = _tf(circle)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    _set_font(run, size=Pt(22), color=BG_DARK, bold=True)
    # Agent name
    _add_text_box(s, x + Inches(0.1), Inches(2.5), Inches(2.5), Inches(0.8),
                  name, size=Pt(16), color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    _add_text_box(s, x + Inches(0.15), Inches(3.5), Inches(2.4), Inches(2.0),
                  desc, size=Pt(12), color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    # Arrow between agents
    if i < 3:
        _add_flow_arrow(s, x + Inches(2.75), Inches(3.5), Inches(0.4), color)

# Post-processing bar at bottom
post_bar = _add_shape(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.8),
                      fill_color=BG_CARD_ALT, border_color=ACCENT_CYAN, border_width=Pt(1.5))
tf = _tf(post_bar)
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "POST-PROCESSING:  UnifiedRiskEngine \u2192 Score Calculation  \u2192  ActionItem Generation  \u2192  WebSocket Broadcast (RISK_UPDATE)"
_set_font(run, size=Pt(13), color=ACCENT_CYAN)

_slide_number(s, 10)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RECON AGENT DEEP DIVE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "RECON AGENT — Deep Dive", "Infrastructure Discovery & Web Crawling")

_add_bullet_card(s, Inches(0.6), Inches(1.2), Inches(5.8), Inches(2.5),
    "Nmap Integration",
    ["\u2022 Quick scan: -F -T4 (top 100 ports)",
     "\u2022 Full scan: -sV -T4 (service detection)",
     "\u2022 Deep scan: -sV -O -A --script=vulners,banner",
     "\u2022 OS fingerprinting, MAC vendor lookup",
     "\u2022 Creates ScanAsset + AssetService records"],
    border=ACCENT_BLUE)

_add_bullet_card(s, Inches(6.7), Inches(1.2), Inches(5.8), Inches(2.5),
    "Playwright Web Crawling",
    ["\u2022 JS-enabled headless browser (Chromium)",
     "\u2022 Discovers links, forms, and input fields",
     "\u2022 Waits for networkidle (SPA-compatible)",
     "\u2022 Fallback to requests if Playwright unavailable",
     "\u2022 Auth credential injection support"],
    border=ACCENT_CYAN)

_add_bullet_card(s, Inches(0.6), Inches(4.0), Inches(5.8), Inches(2.5),
    "Tech Stack Detection",
    ["\u2022 HTTP headers: Server, X-Powered-By",
     "\u2022 Content analysis: React, Vue, Angular",
     "\u2022 CMS detection: WordPress, Drupal",
     "\u2022 Framework fingerprinting",
     "\u2022 Stored in Target.tech_stack (JSON)"],
    border=ACCENT_GREEN)

_add_bullet_card(s, Inches(6.7), Inches(4.0), Inches(5.8), Inches(2.5),
    "Output Schema",
    ["\u2022 endpoints[]: {url, method, parameters}",
     "\u2022 tech_stack: {server, frontend, cms}",
     "\u2022 forms[]: {action, method, inputs[]}",
     "\u2022 assets[]: {ip, ports[], os, hostname}",
     "\u2022 Broadcasts [RECON] events via WebSocket"],
    border=ACCENT_ORANGE)

_slide_number(s, 11)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ATTACK AGENT DEEP DIVE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "ATTACK AGENT — Deep Dive", "Vulnerability Discovery & Payload Testing")

# Nuclei template mapping
_add_bullet_card(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(2.5),
    "Service \u2192 Nuclei Template Mapping",
    ["\u2022 http/https \u2192 cve, exposures, misconfiguration, takeovers",
     "\u2022 ftp/ssh \u2192 default-logins, misconfiguration",
     "\u2022 mysql/postgresql \u2192 default-logins, misconfiguration",
     "\u2022 redis/smb \u2192 cve, default-logins, misconfiguration",
     "\u2022 Timeout: 600s per target, JSONL output parsing"],
    border=ACCENT_RED)

# Payload library
_add_bullet_card(s, Inches(6.9), Inches(1.2), Inches(5.8), Inches(2.5),
    "Built-in Payload Library",
    ["\u2022 SQLi: ' OR '1'='1, UNION SELECT NULL, DROP TABLE",
     "\u2022 XSS: <script>alert(1)</script>, onerror handlers",
     "\u2022 BOLA: ../../../etc/passwd, /api/users/999999",
     "\u2022 SSRF: http://127.0.0.1, file:///etc/passwd",
     "\u2022 Concurrent execution via asyncio.Semaphore"],
    border=ACCENT_ORANGE)

# Lab heuristics
_add_bullet_card(s, Inches(0.6), Inches(4.0), Inches(6.0), Inches(2.8),
    "Lab-Specific Heuristics",
    ["\u2022 Port 6379 (Redis no-auth) \u2192 CRITICAL, confidence 1.0",
     "\u2022 Port 3000 (Juice Shop) \u2192 HIGH, confidence 0.9",
     "\u2022 Port 445 (SMB weak creds) \u2192 CRITICAL, confidence 0.95",
     "\u2022 Port 5432 (PostgreSQL exposed) \u2192 HIGH, confidence 0.85",
     "\u2022 Combines deterministic rules + Nuclei findings"],
    border=ACCENT_RED)

# Concurrency model
_add_bullet_card(s, Inches(6.9), Inches(4.0), Inches(5.8), Inches(2.8),
    "Concurrency Model",
    ["\u2022 asyncio.gather() for parallel payload dispatch",
     "\u2022 httpx.AsyncClient with 10s timeout",
     "\u2022 Up to 50 concurrent endpoint tests",
     "\u2022 Response analysis: status codes + content matching",
     "\u2022 Creates Vulnerability records with evidence JSON"],
    border=ACCENT_BLUE)

_slide_number(s, 12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — VALIDATION AGENT DEEP DIVE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "VALIDATION AGENT — Deep Dive", "False Positive Filtering & Confidence Scoring")

# Two-column layout
_add_bullet_card(s, Inches(0.6), Inches(1.2), Inches(5.8), Inches(3.0),
    "Deterministic Filtering (Stage 1)",
    ["\u2022 Confidence threshold: \u2265 0.6 required",
     "\u2022 Evidence completeness check",
     "\u2022 Cross-reference against known patterns",
     "\u2022 Reject findings without reproducible evidence",
     "\u2022 Status update: OPEN \u2192 FALSE_POSITIVE if rejected",
     "\u2022 100% reproducible \u2014 same input \u2192 same output"],
    border=ACCENT_ORANGE)

_add_bullet_card(s, Inches(6.7), Inches(1.2), Inches(5.8), Inches(3.0),
    "LLM Re-Validation (Stage 2 \u2014 Optional)",
    ["\u2022 Triggered only for findings passing Stage 1",
     "\u2022 Google Gemini 2.0 Flash model",
     "\u2022 Prompt: vulnerability type, URL, evidence, confidence",
     "\u2022 Response: VERDICT (REAL/FALSE_POSITIVE)",
     "\u2022 New confidence score assigned by LLM",
     "\u2022 Non-deterministic \u2014 enhances but does not override"],
    border=ACCENT_CYAN)

# Visual: Funnel
_add_text_box(s, Inches(0.6), Inches(4.5), Inches(12.0), Inches(0.5),
              "VALIDATION FUNNEL", size=Pt(16), color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

funnel_stages = [
    ("ALL FINDINGS", Inches(11.0), ACCENT_RED),
    ("Confidence \u2265 0.6", Inches(8.5), ACCENT_ORANGE),
    ("LLM Verified", Inches(6.0), ACCENT_GREEN),
]

for i, (label, width, color) in enumerate(funnel_stages):
    y = Inches(5.1) + Inches(i * 0.7)
    x = (SLIDE_W - width) / 2
    bar = _add_shape(s, x, y, width, Inches(0.55), fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    tf = _tf(bar)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    _set_font(run, size=Pt(13), color=color, bold=True)

_slide_number(s, 13)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — REPORTING AGENT DEEP DIVE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "REPORTING AGENT — Deep Dive", "AI-Generated Advisory & PoC Scripts")

_add_bullet_card(s, Inches(0.6), Inches(1.2), Inches(5.8), Inches(2.8),
    "Executive Summary Generation",
    ["\u2022 Powered by Google Gemini 2.0 Flash",
     "\u2022 Target audience: Non-technical CEO/CTO",
     "\u2022 3-4 sentences covering security posture",
     "\u2022 Most concerning finding highlighted",
     "\u2022 Immediate next steps recommended",
     "\u2022 Simple language, no jargon policy"],
    border=ACCENT_GREEN)

_add_bullet_card(s, Inches(6.7), Inches(1.2), Inches(5.8), Inches(2.8),
    "Proof-of-Concept Generation",
    ["\u2022 Auto-generated Python scripts for Critical/High",
     "\u2022 Templates: SQLi, XSS, BOLA exploit scripts",
     "\u2022 Includes target URL and payload injection",
     "\u2022 Responsible disclosure format",
     "\u2022 Facilitates remediation verification"],
    border=ACCENT_RED)

_add_bullet_card(s, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.5),
    "Remediation Intelligence",
    ["\u2022 Vulnerability-specific fix recommendations",
     "\u2022 Configuration hardening suggestions",
     "\u2022 Mapped to OWASP best practices",
     "\u2022 Stored in Vulnerability.remediation field",
     "\u2022 Drives ActionItem generation downstream"],
    border=ACCENT_BLUE)

_add_bullet_card(s, Inches(6.7), Inches(4.3), Inches(5.8), Inches(2.5),
    "Report Assembly",
    ["\u2022 Full Markdown report with sections",
     "\u2022 Risk score summary + vulnerability table",
     "\u2022 Asset inventory with port analysis",
     "\u2022 Stored in Scan.agent_thoughts JSON",
     "\u2022 Exportable via Reports dashboard component"],
    border=ACCENT_ORANGE)

_slide_number(s, 14)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — UNIFIED RISK ENGINE OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "UNIFIED RISK ENGINE", "Deterministic Risk Quantification Framework")

# Formula card
formula_card = _add_shape(s, Inches(0.6), Inches(1.2), Inches(12.0), Inches(1.3),
                          fill_color=BG_CARD_ALT, border_color=ACCENT_CYAN, border_width=Pt(2))
tf = _tf(formula_card)
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "RISK SCORE  =  min(100,  PENALTY  \u00D7  ASSET_MULTIPLIER  \u00D7  EXPOSURE_MULTIPLIER)"
_set_font(run, size=Pt(22), color=ACCENT_CYAN, bold=True)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "where  PENALTY = \u03A3(severity_weight \u00D7 confidence) + \u03A3(port_risk_penalty)"
_set_font(run2, size=Pt(14), color=TEXT_LIGHT)

# Three weight tables
_add_bullet_card(s, Inches(0.6), Inches(2.8), Inches(3.7), Inches(3.5),
    "Severity Weights",
    ["CRITICAL  \u2192  25 points",
     "HIGH         \u2192  15 points",
     "MEDIUM    \u2192   7  points",
     "LOW          \u2192   2  points",
     "INFO          \u2192   0  points",
     "",
     "Weighted by confidence_score (0\u20131.0)"],
    title_color=ACCENT_RED, border=ACCENT_RED)

_add_bullet_card(s, Inches(4.6), Inches(2.8), Inches(3.7), Inches(3.5),
    "Port Risk Penalties",
    ["FTP (21)            \u2192  15 pts",
     "Telnet (23)        \u2192  20 pts",
     "SMB (445)         \u2192  20 pts",
     "RDP (3389)       \u2192  15 pts",
     "Redis (6379)     \u2192  10 pts",
     "PostgreSQL (5432)  \u2192  10 pts",
     "MySQL (3306)    \u2192  10 pts"],
    title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

_add_bullet_card(s, Inches(8.6), Inches(2.8), Inches(3.7), Inches(3.5),
    "Asset & Exposure Multipliers",
    ["CRITICAL asset   \u2192  1.5\u00D7",
     "HIGH asset         \u2192  1.2\u00D7",
     "MEDIUM asset    \u2192  1.0\u00D7",
     "LOW asset          \u2192  0.8\u00D7",
     "",
     "Public IP            \u2192  1.0\u00D7",
     "Private (RFC1918) \u2192  0.6\u00D7"],
    title_color=ACCENT_GREEN, border=ACCENT_GREEN)

_slide_number(s, 15)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — RISK ENGINE: WORKED EXAMPLE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "RISK SCORING — Worked Example", "Deterministic Calculation Walkthrough")

# Scenario card
scenario = _add_shape(s, Inches(0.6), Inches(1.2), Inches(12.0), Inches(1.5),
                      fill_color=BG_CARD, border_color=ACCENT_BLUE, border_width=Pt(1.5))
tf = _tf(scenario)
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "SCENARIO:  Lab Redis Cache (10.10.30.20)"
_set_font(run, size=Pt(18), color=ACCENT_CYAN, bold=True)
for line in [
    "Asset Criticality: MEDIUM (1.0\u00D7)  |  Network: Private (0.6\u00D7)  |  Port 6379 open (Redis, no auth)",
    "Finding 1: CRITICAL \u2014 Unprotected Redis Database (confidence: 1.0)  |  Finding 2: HIGH \u2014 Service Exposure (confidence: 0.8)"
]:
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = line
    _set_font(run2, size=Pt(12), color=TEXT_LIGHT)

# Step-by-step calculation
steps = [
    ("Step 1: Vulnerability Penalties", "CRITICAL: 25 \u00D7 1.0 = 25.0  |  HIGH: 15 \u00D7 0.8 = 12.0", ACCENT_RED),
    ("Step 2: Port Risk Penalty", "Redis (6379): +10 points", ACCENT_ORANGE),
    ("Step 3: Total Penalty", "\u03A3 = 25.0 + 12.0 + 10.0 = 47.0", ACCENT_BLUE),
    ("Step 4: Apply Multipliers", "47.0 \u00D7 1.0 (MEDIUM) \u00D7 0.6 (Private) = 28.2", ACCENT_GREEN),
    ("Step 5: Final Score", "RISK SCORE = min(100, 28.2) = 28.2 / 100", ACCENT_CYAN),
]

for i, (step, calc, color) in enumerate(steps):
    y = Inches(3.0) + Inches(i * 0.85)
    card = _add_shape(s, Inches(0.6), y, Inches(12.0), Inches(0.7),
                      fill_color=BG_CARD, border_color=color, border_width=Pt(1))
    _add_text_box(s, Inches(0.9), y + Inches(0.1), Inches(3.5), Inches(0.4),
                  step, size=Pt(13), color=color, bold=True)
    _add_text_box(s, Inches(4.5), y + Inches(0.1), Inches(7.8), Inches(0.4),
                  calc, size=Pt(13), color=TEXT_LIGHT)

_slide_number(s, 16)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — DETERMINISTIC vs NON-DETERMINISTIC
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "DETERMINISTIC vs. NON-DETERMINISTIC LOGIC", "Hybrid Intelligence Architecture")

# Left: Deterministic
det_card = _add_shape(s, Inches(0.6), Inches(1.2), Inches(5.8), Inches(5.5),
                      fill_color=BG_CARD, border_color=ACCENT_GREEN, border_width=Pt(2))
tf = _tf(det_card)
tf.margin_left = Inches(0.25)
tf.margin_top = Inches(0.2)
tf.margin_right = Inches(0.2)

p = tf.paragraphs[0]
run = p.add_run()
run.text = "DETERMINISTIC COMPONENTS"
_set_font(run, size=Pt(18), color=ACCENT_GREEN, bold=True)

det_items = [
    "\u2713 Risk Score Calculation (UnifiedRiskEngine)",
    "    Same vulns + ports + criticality = Same score",
    "\u2713 Nmap Port Scanning (NmapWrapper)",
    "    Reproducible infrastructure discovery",
    "\u2713 Nuclei Template Matching",
    "    Known CVE patterns, no ambiguity",
    "\u2713 Confidence Threshold Filtering (\u22650.6)",
    "    Binary pass/fail gate",
    "\u2713 Health Score Computation",
    "    Fixed deduction rules per severity",
    "\u2713 Action Item Generation",
    "    Rule-based: HIGH/CRITICAL \u2192 REMEDIATION",
    "\u2713 Port Risk Penalty Table",
    "    Static mapping: port \u2192 penalty points",
]

for item in det_items:
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = item
    is_sub = item.startswith("    ")
    _set_font(run2, size=Pt(11) if is_sub else Pt(12),
              color=TEXT_DIM if is_sub else TEXT_LIGHT,
              italic=is_sub)
    p2.space_before = Pt(1)
    p2.space_after = Pt(1)

# Right: Non-Deterministic
ndet_card = _add_shape(s, Inches(6.7), Inches(1.2), Inches(5.8), Inches(5.5),
                       fill_color=BG_CARD, border_color=ACCENT_CYAN, border_width=Pt(2))
tf = _tf(ndet_card)
tf.margin_left = Inches(0.25)
tf.margin_top = Inches(0.2)
tf.margin_right = Inches(0.2)

p = tf.paragraphs[0]
run = p.add_run()
run.text = "NON-DETERMINISTIC COMPONENTS"
_set_font(run, size=Pt(18), color=ACCENT_CYAN, bold=True)

ndet_items = [
    "\u25CB LLM Re-Validation (Gemini 2.0 Flash)",
    "    VERDICT: REAL/FALSE_POSITIVE + confidence",
    "\u25CB Executive Summary Generation",
    "    Natural language for non-technical audience",
    "\u25CB PoC Script Enhancement",
    "    LLM-augmented exploit code generation",
    "\u25CB Asset Intelligence Advisory",
    "    Risk explanation, business impact, remediation",
    "\u25CB SIEM Event Analysis",
    "    AI interpretation of Wazuh/ES alerts",
    "",
    "KEY DESIGN PRINCIPLE:",
    "Non-deterministic outputs ENHANCE but never",
    "OVERRIDE deterministic scoring decisions.",
]

for item in ndet_items:
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = item
    is_sub = item.startswith("    ")
    is_key = "KEY DESIGN" in item or "OVERRIDE" in item or "Non-deterministic" in item
    _set_font(run2, size=Pt(11) if is_sub else Pt(12),
              color=ACCENT_ORANGE if is_key else (TEXT_DIM if is_sub else TEXT_LIGHT),
              bold=is_key, italic=is_sub)
    p2.space_before = Pt(1)
    p2.space_after = Pt(1)

_slide_number(s, 17)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — ASSET CRITICALITY WEIGHTING
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "ASSET CRITICALITY WEIGHTING", "Context-Aware Risk Amplification")

# Explanation
_add_text_box(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.8),
              "The Unified Risk Engine does not treat all assets equally. A CRITICAL vulnerability on a production "
              "database server carries more organizational risk than the same vulnerability on an isolated test workstation. "
              "Asset criticality multipliers amplify or attenuate the raw penalty score.",
              size=Pt(13), color=TEXT_LIGHT)

# Chart: Bar chart showing multiplier impact
chart_data = CategoryChartData()
chart_data.categories = ['CRITICAL\n(1.5\u00D7)', 'HIGH\n(1.2\u00D7)', 'MEDIUM\n(1.0\u00D7)', 'LOW\n(0.8\u00D7)']
chart_data.add_series('Risk Score (Penalty = 50)', (75, 60, 50, 40))
chart_data.add_series('Risk Score (Penalty = 30)', (45, 36, 30, 24))

chart_frame = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(2.2), Inches(7.0), Inches(4.5),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(11)
chart.legend.font.color.rgb = TEXT_LIGHT

plot = chart.plots[0]
plot.gap_width = 80
series0 = plot.series[0]
series1 = plot.series[1]
series0.format.fill.solid()
series0.format.fill.fore_color.rgb = ACCENT_CYAN
series1.format.fill.solid()
series1.format.fill.fore_color.rgb = ACCENT_BLUE

# Data labels
series0.has_data_labels = True
series0.data_labels.font.size = Pt(11)
series0.data_labels.font.color.rgb = TEXT_WHITE
series0.data_labels.number_format = '0'
series1.has_data_labels = True
series1.data_labels.font.size = Pt(11)
series1.data_labels.font.color.rgb = TEXT_WHITE
series1.data_labels.number_format = '0'

# Chart styling
value_axis = chart.value_axis
value_axis.visible = True
value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x2A, 0x30, 0x50)
value_axis.format.line.color.rgb = TEXT_DIM
value_axis.tick_labels.font.size = Pt(10)
value_axis.tick_labels.font.color.rgb = TEXT_DIM
value_axis.maximum_scale = 100

category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(10)
category_axis.tick_labels.font.color.rgb = TEXT_LIGHT
category_axis.format.line.color.rgb = TEXT_DIM

chart.chart_style = 2  # minimal style

# Exposure card
_add_bullet_card(s, Inches(8.0), Inches(2.2), Inches(4.5), Inches(4.5),
    "Network Exposure Factor",
    ["\u2022 Public IP assets: 1.0\u00D7 (full risk)",
     "\u2022 Private (RFC 1918): 0.6\u00D7 (attenuated)",
     "",
     "\u2022 Rationale: Internal assets require",
     "  attacker to first breach perimeter,",
     "  reducing effective exposure.",
     "",
     "\u2022 Combined: Criticality \u00D7 Exposure",
     "  e.g., CRITICAL + Public = 1.5 \u00D7 1.0",
     "  e.g., LOW + Private = 0.8 \u00D7 0.6 = 0.48"],
    title_color=ACCENT_GREEN, border=ACCENT_GREEN)

_slide_number(s, 18)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — HEALTH SCORE COMPUTATION
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "HEALTH SCORE COMPUTATION", "Complementary Metric (100 = Safe)")

formula_card = _add_shape(s, Inches(0.6), Inches(1.2), Inches(12.0), Inches(1.0),
                          fill_color=BG_CARD_ALT, border_color=ACCENT_GREEN, border_width=Pt(2))
tf = _tf(formula_card)
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "HEALTH = max(0,  100  \u2212  \u03A3(vuln_deductions)  \u2212  \u03A3(port_deductions))    |    Capped at 90 if any vulnerabilities exist"
_set_font(run, size=Pt(16), color=ACCENT_GREEN, bold=True)

# Deduction tables
_add_bullet_card(s, Inches(0.6), Inches(2.5), Inches(5.5), Inches(2.5),
    "Vulnerability Deductions",
    ["Each CRITICAL vulnerability:   \u221220 points",
     "Each HIGH vulnerability:           \u221210 points",
     "Each MEDIUM vulnerability:      \u2212 5  points",
     "Each LOW vulnerability:             \u2212 0  points",
     "Cap: score cannot exceed 90 if vulns present"],
    title_color=ACCENT_RED, border=ACCENT_RED)

_add_bullet_card(s, Inches(6.4), Inches(2.5), Inches(6.0), Inches(2.5),
    "High-Risk Port Deductions",
    ["Each open FTP (21):       \u221215 points",
     "Each open Telnet (23):   \u221215 points",
     "Each open SMB (445):   \u221215 points",
     "Each open RDP (3389): \u221215 points",
     "Floor: health score never drops below 0"],
    title_color=ACCENT_ORANGE, border=ACCENT_ORANGE)

# Example
example = _add_shape(s, Inches(0.6), Inches(5.3), Inches(12.0), Inches(1.5),
                     fill_color=BG_CARD, border_color=ACCENT_CYAN, border_width=Pt(1))
tf = _tf(example)
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "EXAMPLE:  2 CRITICAL + 1 HIGH + Open Telnet"
_set_font(run, size=Pt(15), color=ACCENT_CYAN, bold=True)
p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = "Health = 100 \u2212 (2\u00D720) \u2212 (1\u00D710) \u2212 15 = 100 \u2212 40 \u2212 10 \u2212 15 = 35    |    Capped at min(35, 90) = 35 / 100"
_set_font(run2, size=Pt(13), color=TEXT_LIGHT)
p3 = tf.add_paragraph()
run3 = p3.add_run()
run3.text = "Interpretation: Asset is in POOR health \u2014 immediate remediation required."
_set_font(run3, size=Pt(12), color=ACCENT_RED)

_slide_number(s, 19)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — ACTION ITEM GENERATION
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "ACTION ITEM GENERATION", "Automated Remediation Tracking")

_add_bullet_card(s, Inches(0.6), Inches(1.2), Inches(5.8), Inches(3.5),
    "Rule-Based Generation (Deterministic)",
    ["\u2022 Rule 1: Critical/High vulns \u2192 REMEDIATION items",
     "    title = vuln.title or 'Fix {vuln.type}'",
     "    priority = severity.upper()",
     "    status = 'OPEN'",
     "",
     "\u2022 Rule 2: High-risk ports \u2192 CONFIGURATION items",
     "    Ports: 21, 23, 445, 3389, 6379",
     "    title = 'Secure {service} on {ip}'",
     "    description = 'Restrict or VPN access'",
     "",
     "\u2022 Deduplication by title to avoid duplicate items"],
    border=ACCENT_BLUE)

_add_bullet_card(s, Inches(6.7), Inches(1.2), Inches(5.8), Inches(3.5),
    "ActionItem Lifecycle",
    ["\u2022 OPEN \u2192 IN_PROGRESS \u2192 RESOLVED",
     "\u2022 OPEN \u2192 ACCEPTED (risk accepted)",
     "",
     "Item Types:",
     "\u2022 REMEDIATION \u2014 Fix a vulnerability",
     "\u2022 CONFIGURATION \u2014 Harden infrastructure",
     "\u2022 new_device \u2014 New device discovered",
     "\u2022 alert \u2014 SIEM/Wazuh notification",
     "",
     "Dashboard Integration:",
     "\u2022 GET /api/v1/dashboard/actions"],
    border=ACCENT_GREEN)

# Visual flow
flow_card = _add_shape(s, Inches(0.6), Inches(5.0), Inches(12.0), Inches(1.5),
                       fill_color=BG_CARD_ALT, border_color=ACCENT_CYAN, border_width=Pt(1))

flow_items = [
    ("Scan Complete", Inches(0.8), ACCENT_BLUE),
    ("Risk Engine", Inches(3.3), ACCENT_ORANGE),
    ("Action Items", Inches(5.8), ACCENT_GREEN),
    ("Dashboard", Inches(8.3), ACCENT_CYAN),
    ("Resolved", Inches(10.8), ACCENT_GREEN),
]

for label, x, color in flow_items:
    box = _add_shape(s, x, Inches(5.3), Inches(2.0), Inches(0.8),
                     fill_color=BG_CARD, border_color=color, border_width=Pt(1))
    tf = _tf(box)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    _set_font(run, size=Pt(11), color=color, bold=True)

for x in [Inches(2.85), Inches(5.35), Inches(7.85), Inches(10.35)]:
    _add_flow_arrow(s, x, Inches(5.55), Inches(0.4), ACCENT_CYAN)

_slide_number(s, 20)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — CELERY TASK PIPELINE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "CELERY TASK PIPELINE", "Asynchronous Scan Execution & Scheduling")

# Pipeline phases
phases = [
    ("Phase 1: NMAP", "Nmap scan \u2192 ScanAsset + AssetService records \u2192 Vulnerability (Service Exposure)", ACCENT_BLUE),
    ("Phase 2: AI INTELLIGENCE", "Gemini batch_analyze \u2192 ai_insight JSON per asset (risk explanation, business impact)", ACCENT_CYAN),
    ("Phase 3: NUCLEI", "Template-based vulnerability detection \u2192 Vulnerability records with CVE info", ACCENT_RED),
    ("Phase 4: RISK ENGINE", "UnifiedRiskEngine.update_scan_risk() \u2192 risk_score + health_score + action_items", ACCENT_GREEN),
    ("Phase 5: ASSET MONITOR", "AssetMonitor.process_scan_results() \u2192 Detect new devices, update NetworkAsset", ACCENT_ORANGE),
    ("Phase 6: EVENT PUBLISH", "Redis Pub/Sub \u2192 WebSocket broadcast \u2192 Frontend RISK_UPDATE event", ACCENT_CYAN),
]

for i, (phase, desc, color) in enumerate(phases):
    y = Inches(1.2) + Inches(i * 0.9)
    card = _add_shape(s, Inches(0.6), y, Inches(12.0), Inches(0.75),
                      fill_color=BG_CARD, border_color=color, border_width=Pt(1))
    badge = _add_shape(s, Inches(0.8), y + Inches(0.1), Inches(2.8), Inches(0.5),
                       fill_color=color)
    tf = _tf(badge)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = phase
    _set_font(run, size=Pt(11), color=BG_DARK, bold=True)

    _add_text_box(s, Inches(3.8), y + Inches(0.15), Inches(8.5), Inches(0.5),
                  desc, size=Pt(12), color=TEXT_LIGHT)

# Celery Beat scheduling
beat_card = _add_shape(s, Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.7),
                       fill_color=BG_CARD_ALT, border_color=ACCENT_ORANGE, border_width=Pt(1))
tf = _tf(beat_card)
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "CELERY BEAT:  Hourly automatic scans via crontab(minute=0)  |  max_retries=2  |  retry_delay=30s  |  @celery_app.task(bind=True)"
_set_font(run, size=Pt(12), color=ACCENT_ORANGE)

_slide_number(s, 21)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — SIEM / SOAR INTEGRATION
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "SIEM / SOAR INTEGRATION", "Wazuh + Elasticsearch + n8n")

_add_bullet_card(s, Inches(0.6), Inches(1.2), Inches(3.8), Inches(3.0),
    "Wazuh Manager",
    ["\u2022 Security event management",
     "\u2022 Agent monitoring (port 1514/1515)",
     "\u2022 REST API on port 55000",
     "\u2022 Active response triggers",
     "\u2022 JWT authentication",
     "\u2022 Integrated alert forwarding"],
    border=ACCENT_BLUE)

_add_bullet_card(s, Inches(4.6), Inches(1.2), Inches(3.8), Inches(3.0),
    "Elasticsearch + Kibana",
    ["\u2022 Log aggregation (port 9200)",
     "\u2022 Index: sme-lab-events-*",
     "\u2022 Index: wazuh-alerts-*",
     "\u2022 ECS format compliance",
     "\u2022 Kibana dashboards (5601)",
     "\u2022 Bulk indexing via /_bulk API"],
    border=ACCENT_GREEN)

_add_bullet_card(s, Inches(8.6), Inches(1.2), Inches(3.8), Inches(3.0),
    "n8n SOAR",
    ["\u2022 Workflow automation (port 5678)",
     "\u2022 Webhook-triggered responses",
     "\u2022 Connects to Orchestration Security Center API",
     "\u2022 Automated incident workflows",
     "\u2022 Slack/Email notifications",
     "\u2022 Remediation playbooks"],
    border=ACCENT_ORANGE)

# Log flow visualization
_add_text_box(s, Inches(0.6), Inches(4.5), Inches(12.0), Inches(0.4),
              "LOG FLOW ARCHITECTURE", size=Pt(16), color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

flow_items = [
    ("Traffic Gen", Inches(0.3), ACCENT_RED),
    ("Log Shipper", Inches(2.6), ACCENT_ORANGE),
    ("Elasticsearch", Inches(4.9), ACCENT_GREEN),
    ("Wazuh", Inches(7.2), ACCENT_BLUE),
    ("Dashboard", Inches(9.5), ACCENT_CYAN),
    ("n8n SOAR", Inches(11.5), ACCENT_ORANGE),
]

for label, x, color in flow_items:
    box = _add_shape(s, x, Inches(5.0), Inches(2.0), Inches(1.0),
                     fill_color=BG_CARD, border_color=color, border_width=Pt(1))
    tf = _tf(box)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    _set_font(run, size=Pt(12), color=color, bold=True)

for x in [Inches(2.35), Inches(4.65), Inches(6.95), Inches(9.25), Inches(11.25)]:
    _add_flow_arrow(s, x, Inches(5.25), Inches(0.3), ACCENT_CYAN)

# Wazuh rules
_add_bullet_card(s, Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.9),
    "Wazuh Rule Mapping Examples",
    ["\u2022 Auth Failure \u2192 Rule 5710 (Level 5)    \u2022 Port Scan \u2192 Rule 510 (Level 6)    \u2022 Brute Force \u2192 Rule 5712 (Level 10)    \u2022 DB Error \u2192 Rule 50100 (Level 7)"],
    border=ACCENT_BLUE)

_slide_number(s, 22)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — LIVING LAB ENVIRONMENT
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "LIVING LAB ENVIRONMENT", "4 Subnets, 11 Vulnerable Containers")

subnets = [
    ("DMZ (10.10.10.0/24)", ACCENT_RED, [
        ("Juice Shop", "10.10.10.10:3000", "OWASP Top 10 vulns"),
        ("API Gateway", "10.10.10.20:8081", "Nginx exposed headers"),
        ("DNS Server", "10.10.10.30:53", "Zone transfer vuln"),
    ]),
    ("Corporate (10.10.20.0/24)", ACCENT_ORANGE, [
        ("File Server", "10.10.20.10:445", "SMB weak creds"),
        ("Mail Server", "10.10.20.20", "Plaintext SMTP/POP3"),
        ("Workstation", "10.10.20.40", "Internal exposure"),
    ]),
    ("Data (10.10.30.0/24)", ACCENT_CYAN, [
        ("PostgreSQL", "10.10.30.10:5432", "password123"),
        ("Redis Cache", "10.10.30.20:6380", "No auth, no protected-mode"),
    ]),
    ("Management (10.10.40.0/24)", ACCENT_GREEN, [
        ("Traffic Gen", "10.10.40.x", "7 traffic generators"),
        ("Log Shipper", "10.10.40.x", "ECS + Wazuh format"),
    ]),
]

for i, (subnet_name, color, hosts) in enumerate(subnets):
    x = Inches(0.4) + Inches(i * 3.2)
    # Subnet header
    header = _add_shape(s, x, Inches(1.2), Inches(3.0), Inches(0.5),
                        fill_color=color)
    tf = _tf(header)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = subnet_name
    _set_font(run, size=Pt(12), color=BG_DARK, bold=True)

    # Hosts
    for j, (name, addr, vuln) in enumerate(hosts):
        y = Inches(1.85) + Inches(j * 1.5)
        card = _add_shape(s, x, y, Inches(3.0), Inches(1.3),
                          fill_color=BG_CARD, border_color=color, border_width=Pt(1))
        tf = _tf(card)
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.08)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = name
        _set_font(run, size=Pt(13), color=TEXT_WHITE, bold=True)
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = addr
        _set_font(run2, size=Pt(10), color=ACCENT_CYAN)
        p3 = tf.add_paragraph()
        run3 = p3.add_run()
        run3.text = vuln
        _set_font(run3, size=Pt(10), color=TEXT_DIM)

_slide_number(s, 23)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — TRAFFIC GENERATION & LOG SHIPPING
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "TRAFFIC GENERATION & LOG SHIPPING", "Simulated Enterprise Activity")

_add_bullet_card(s, Inches(0.6), Inches(1.2), Inches(5.8), Inches(2.8),
    "7 Traffic Generators (Threaded)",
    ["\u2022 HTTP: Employee browsing + API calls",
     "\u2022 DNS: Internal + external domain queries",
     "\u2022 SMB: File server access (admin/admin123)",
     "\u2022 Database: PostgreSQL queries from webserver",
     "\u2022 Redis: GET/SET cache operations",
     "\u2022 Email: SMTP/POP3/IMAP simulated mail",
     "\u2022 Suspicious: Port scans, brute force, exfiltration"],
    border=ACCENT_RED)

_add_bullet_card(s, Inches(6.7), Inches(1.2), Inches(5.8), Inches(2.8),
    "Log Shipper Pipeline",
    ["\u2022 Tail-follow /var/log/traffic/traffic.log",
     "\u2022 Transform to ECS format (Elastic Common Schema)",
     "\u2022 Transform to Wazuh alert format",
     "\u2022 Buffer in memory (size \u2265 50 or 10s interval)",
     "\u2022 Bulk-index to Elasticsearch via /_bulk API",
     "\u2022 Creates sme-lab-events-* and wazuh-alerts-*",
     "\u2022 Index templates with proper field mappings"],
    border=ACCENT_GREEN)

# Intensity levels chart
chart_data = CategoryChartData()
chart_data.categories = ['HTTP', 'DNS', 'SMB', 'DB', 'Redis', 'Email', 'Suspicious']
chart_data.add_series('Low (seconds)', (30, 45, 60, 40, 20, 120, 300))
chart_data.add_series('Medium', (10, 15, 30, 20, 10, 60, 120))
chart_data.add_series('High', (3, 5, 10, 8, 5, 30, 60))

chart_frame = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(4.3), Inches(12.0), Inches(3.0),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.font.size = Pt(10)
chart.legend.font.color.rgb = TEXT_LIGHT
chart.legend.include_in_layout = False

plot = chart.plots[0]
plot.gap_width = 60
colors = [ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED]
for i, color in enumerate(colors):
    series = plot.series[i]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color

value_axis = chart.value_axis
value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x2A, 0x30, 0x50)
value_axis.tick_labels.font.size = Pt(9)
value_axis.tick_labels.font.color.rgb = TEXT_DIM
value_axis.format.line.color.rgb = TEXT_DIM
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(9)
category_axis.tick_labels.font.color.rgb = TEXT_LIGHT
category_axis.format.line.color.rgb = TEXT_DIM

_slide_number(s, 24)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 25 — FRONTEND DASHBOARD
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "DASHBOARD & REAL-TIME VISUALIZATION", "React/Vite SPA with WebSocket Updates")

components = [
    ("RiskScore Gauge", "0\u2013100 gauge with color coding\nGreen (\u226570) / Orange (40\u201370) / Red (<40)", ACCENT_CYAN),
    ("Scan Pipeline", "Visual step progress: Queued \u2192 Nmap\n\u2192 Nuclei \u2192 Risk Engine \u2192 AI", ACCENT_BLUE),
    ("Agent Log Viewer", "Console-like terminal stream\nReal-time fetch every 2 seconds", ACCENT_GREEN),
    ("Action Center", "Open remediation items\nStatus: OPEN \u2192 RESOLVED", ACCENT_ORANGE),
    ("Activity Feed", "Recent security events\nPrioritized by severity", ACCENT_RED),
    ("Network Topology", "Graph visualization of discovered\nassets and connections", ACCENT_BLUE),
    ("Risk Heatmap", "Asset risk matrix with\ncolor-coded severity grid", ACCENT_RED),
    ("Lab Environment", "Lab container status\nTelemetry and controls", ACCENT_GREEN),
]

for i, (name, desc, color) in enumerate(components):
    col = i % 4
    row = i // 4
    x = Inches(0.4) + Inches(col * 3.2)
    y = Inches(1.2) + Inches(row * 3.0)
    card = _add_shape(s, x, y, Inches(3.0), Inches(2.5),
                      fill_color=BG_CARD, border_color=color, border_width=Pt(1.5))
    tf = _tf(card)
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name
    _set_font(run, size=Pt(14), color=color, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    run2 = p2.add_run()
    run2.text = desc
    _set_font(run2, size=Pt(11), color=TEXT_LIGHT)

_slide_number(s, 25)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 26 — REAL-TIME EVENT ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "REAL-TIME EVENT ARCHITECTURE", "Redis Pub/Sub \u2192 WebSocket \u2192 React Context")

# Three-stage flow
stages = [
    ("BACKEND\n(Event Publisher)", "Redis Pub/Sub channel\npublish('RISK_UPDATE', payload)\nTriggered after Risk Engine", ACCENT_BLUE, Inches(0.4)),
    ("GATEWAY\n(FastAPI WebSocket)", "Endpoint: /ws/logs\nRedis listener with backoff\n2s \u2192 4s \u2192 8s \u2192 32s retry", ACCENT_CYAN, Inches(4.5)),
    ("FRONTEND\n(RealTimeContext)", "WebSocket.onmessage\ndispatch({type, payload})\nAll components re-render", ACCENT_GREEN, Inches(8.6)),
]

for label, desc, color, x in stages:
    card = _add_shape(s, x, Inches(1.3), Inches(3.8), Inches(3.0),
                      fill_color=BG_CARD, border_color=color, border_width=Pt(2))
    tf = _tf(card)
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    _set_font(run, size=Pt(16), color=color, bold=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = desc
    _set_font(run2, size=Pt(12), color=TEXT_LIGHT)

_add_flow_arrow(s, Inches(4.2), Inches(2.6), Inches(0.4), ACCENT_CYAN)
_add_flow_arrow(s, Inches(8.3), Inches(2.6), Inches(0.4), ACCENT_GREEN)

# Event types
events_card = _add_shape(s, Inches(0.4), Inches(4.6), Inches(12.2), Inches(2.5),
                         fill_color=BG_CARD_ALT, border_color=ACCENT_CYAN, border_width=Pt(1))
tf = _tf(events_card)
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "WebSocket Event Types"
_set_font(run, size=Pt(16), color=ACCENT_CYAN, bold=True)

ws_events = [
    "RISK_UPDATE  \u2192  {scan_id, overall_score, health_score, vuln_count}  \u2192  Updates RiskScore gauge + KPIs",
    "[RECON] start_recon  \u2192  {target_url, scan_type}  \u2192  Agent Log Viewer + Pipeline Step",
    "[ATTACK] test_complete  \u2192  {findings_count, tested_endpoints}  \u2192  Activity Feed",
    "[VALIDATION] complete  \u2192  {validated_count, false_positives}  \u2192  Status indicators",
    "[REPORTING] complete  \u2192  {report_summary, poc_count}  \u2192  Reports panel update",
]

for evt in ws_events:
    p2 = tf.add_paragraph()
    p2.space_before = Pt(3)
    run2 = p2.add_run()
    run2.text = f"\u2022  {evt}"
    _set_font(run2, size=Pt(11), color=TEXT_LIGHT)

_slide_number(s, 26)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 27 — SECURITY CAPABILITIES MATRIX
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "SECURITY SCANNING CAPABILITIES MATRIX", "Tool Coverage & Intelligence Layers")

capabilities = [
    ("Port Scanning", "Nmap", "ReconAgent", "\u2713", "\u2014", ACCENT_GREEN),
    ("Service Detection", "Nmap -sV", "ReconAgent", "\u2713", "\u2014", ACCENT_GREEN),
    ("OS Detection", "Nmap -O", "ReconAgent", "\u2713", "\u2014", ACCENT_GREEN),
    ("Web App Crawling", "Playwright", "ReconAgent", "\u2713", "\u2014", ACCENT_GREEN),
    ("Template Vulns", "Nuclei", "AttackAgent", "\u2713", "\u2014", ACCENT_GREEN),
    ("Payload Injection", "AttackAgent", "AttackAgent", "\u2713", "\u2014", ACCENT_GREEN),
    ("FP Filtering", "Validation", "ValidationAgent", "\u2713", "Gemini", ACCENT_CYAN),
    ("Risk Scoring", "UnifiedRiskEngine", "All", "\u2713", "\u2014", ACCENT_GREEN),
    ("PoC Generation", "Templates+LLM", "ReportingAgent", "Template", "Gemini", ACCENT_CYAN),
    ("Exec Summary", "LLM", "ReportingAgent", "\u2014", "Gemini", ACCENT_CYAN),
    ("Asset Intel", "Gemini", "IntelligenceAgent", "\u2014", "Gemini", ACCENT_CYAN),
]

# Headers
headers = ["Capability", "Tool", "Agent", "Deterministic", "AI-Enhanced"]
for j, h in enumerate(headers):
    x = Inches(0.6) + Inches(j * 2.5)
    _add_text_box(s, x, Inches(1.1), Inches(2.3), Inches(0.3),
                  h, size=Pt(11), color=ACCENT_CYAN, bold=True)

for i, (cap, tool, agent, det, ai, clr) in enumerate(capabilities):
    y = Inches(1.5) + Inches(i * 0.48)
    row_bg = BG_CARD if i % 2 == 0 else BG_CARD_ALT
    _add_shape(s, Inches(0.6), y, Inches(12.0), Inches(0.43), fill_color=row_bg)
    vals = [cap, tool, agent, det, ai]
    for j, val in enumerate(vals):
        x = Inches(0.7) + Inches(j * 2.5)
        c = clr if j >= 3 and val not in ["\u2014", ""] else TEXT_LIGHT
        _add_text_box(s, x, y + Inches(0.05), Inches(2.3), Inches(0.3),
                      val, size=Pt(10), color=c)

_slide_number(s, 27)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 28 — API ENDPOINTS OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "API ENDPOINTS OVERVIEW", "FastAPI REST + WebSocket Gateway")

endpoint_groups = [
    ("Scans", ACCENT_CYAN, [
        "POST /scans         Create sync scan (Celery)",
        "POST /scans/ai      Create AI scan (Orchestrator)",
        "GET  /scans          List scans (filter by status)",
        "GET  /scans/{id}    Detailed scan + vulns + assets",
    ]),
    ("Dashboard", ACCENT_GREEN, [
        "GET /dashboard/risk-overview    KPI metrics",
        "GET /dashboard/kpi-snapshot     Scores + counts",
        "GET /dashboard/actions              Action items",
        "POST /dashboard/refresh-risk    Force recalc",
    ]),
    ("Network", ACCENT_BLUE, [
        "GET /network/assets           All network assets",
        "GET /network/assets/new     Last 24h discoveries",
        "GET /network/activity          Recent events",
    ]),
    ("SIEM & Lab", ACCENT_ORANGE, [
        "GET /siem/alerts    Elasticsearch alerts",
        "GET /siem/health   ES connectivity check",
        "GET /lab/status      Container status",
        "POST /lab/seed     Register lab targets",
    ]),
]

for i, (group, color, endpoints) in enumerate(endpoint_groups):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + Inches(col * 6.4)
    y = Inches(1.2) + Inches(row * 3.0)
    _add_bullet_card(s, x, y, Inches(6.1), Inches(2.7),
                     group, endpoints, title_color=color, border=color)

_slide_number(s, 28)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 29 — VULNERABILITY SEVERITY DISTRIBUTION (Chart)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "VULNERABILITY ANALYSIS — Lab Results", "Severity Distribution & Coverage")

# Pie chart: Severity distribution
chart_data = CategoryChartData()
chart_data.categories = ['Critical', 'High', 'Medium', 'Low', 'Info']
chart_data.add_series('Count', (4, 8, 12, 6, 3))

chart_frame = s.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.5),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.font.size = Pt(11)
chart.legend.font.color.rgb = TEXT_LIGHT
chart.legend.include_in_layout = False

plot = chart.plots[0]
series = plot.series[0]
colors_pie = [ACCENT_RED, ACCENT_ORANGE, RGBColor(0xFF, 0xD7, 0x00), ACCENT_GREEN, ACCENT_BLUE]
for i, color in enumerate(colors_pie):
    pt = series.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = color

series.has_data_labels = True
series.data_labels.font.size = Pt(11)
series.data_labels.font.color.rgb = TEXT_WHITE
series.data_labels.number_format = '0'
series.data_labels.show_category_name = True
series.data_labels.show_value = True

# Findings summary
_add_bullet_card(s, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5),
    "Lab Scan Findings Summary",
    ["\u2022 33 total vulnerabilities discovered",
     "\u2022 4 CRITICAL: Redis no-auth, SMB weak creds,",
     "  PostgreSQL exposed, DNS zone transfer",
     "\u2022 8 HIGH: Juice Shop vulns, default logins,",
     "  service exposures on sensitive ports",
     "\u2022 12 MEDIUM: HTTP misconfiguration,",
     "  directory listing, exposed headers",
     "\u2022 6 LOW: Information disclosure",
     "\u2022 3 INFO: Version banners",
     "",
     "\u2022 False positive rate: <5% after validation",
     "\u2022 Avg scan time: 3-5 min (quick mode)",
     "\u2022 100% deterministic risk scoring"],
    border=ACCENT_CYAN)

_slide_number(s, 29)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 30 — NOVEL CONTRIBUTIONS
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "NOVEL CONTRIBUTIONS", "Research Impact & Innovation")

contributions = [
    ("1", "Unified Risk Engine",
     "Deterministic scoring algorithm combining CVSS-like severity weights, asset criticality multipliers, and network exposure factors. Reproducible for compliance: same inputs always produce the same risk score.",
     ACCENT_CYAN),
    ("2", "Multi-Agent Autonomous Pipeline",
     "4-stage orchestrated DAST workflow (Recon \u2192 Attack \u2192 Validation \u2192 Reporting) that operates without human intervention, reducing mean-time-to-detection from days to minutes.",
     ACCENT_BLUE),
    ("3", "SME-Focused Simplification",
     "AI-powered translation of technical CVE data and scan results into non-technical business language via Google Gemini, enabling CEO/CTO-level decision-making without security expertise.",
     ACCENT_GREEN),
    ("4", "Living Lab Simulation",
     "4-subnet network with 11 vulnerable containers modeling a realistic SME topology (DMZ, Corporate, Data, Management) with traffic generation and SIEM-compatible log shipping.",
     ACCENT_ORANGE),
    ("5", "Hybrid Intelligence Architecture",
     "Clean separation between deterministic (scoring, filtering) and non-deterministic (LLM advisory) components. AI enhances but never overrides deterministic security decisions.",
     ACCENT_RED),
]

for i, (num, title, desc, color) in enumerate(contributions):
    y = Inches(1.1) + Inches(i * 1.2)
    # Number badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.1), Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = _tf(badge)
    tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    _set_font(run, size=Pt(16), color=BG_DARK, bold=True)

    _add_text_box(s, Inches(1.3), y, Inches(2.5), Inches(0.4),
                  title, size=Pt(15), color=color, bold=True)
    _add_text_box(s, Inches(1.3), y + Inches(0.4), Inches(11.0), Inches(0.7),
                  desc, size=Pt(11), color=TEXT_LIGHT)

_slide_number(s, 30)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 31 — QUANTITATIVE METRICS
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "QUANTITATIVE METRICS", "Platform Performance & Coverage")

# KPI boxes row 1
kpis = [
    ("Microservices", "11", ACCENT_CYAN),
    ("Agent Types", "4", ACCENT_BLUE),
    ("Integration\nPoints", "5+", ACCENT_GREEN),
    ("Lab Containers", "11", ACCENT_ORANGE),
    ("Network\nSubnets", "4", ACCENT_RED),
]

for i, (label, value, color) in enumerate(kpis):
    x = Inches(0.4) + Inches(i * 2.6)
    _add_kpi_box(s, x, Inches(1.3), Inches(2.3), Inches(1.5), label, value, color)

# KPI boxes row 2
kpis2 = [
    ("Concurrent\nEndpoints", "50", ACCENT_CYAN),
    ("Quick Scan", "2-5\nmin", ACCENT_GREEN),
    ("Deep Scan", "15+\nmin", ACCENT_ORANGE),
    ("FP Rate", "<5%", ACCENT_GREEN),
    ("DB Tables", "9", ACCENT_BLUE),
]

for i, (label, value, color) in enumerate(kpis2):
    x = Inches(0.4) + Inches(i * 2.6)
    _add_kpi_box(s, x, Inches(3.2), Inches(2.3), Inches(1.5), label, value, color)

# Technical highlights
_add_bullet_card(s, Inches(0.4), Inches(5.0), Inches(12.2), Inches(2.0),
    "Platform Highlights",
    ["\u2022 Risk scoring is 100% deterministic \u2014 same inputs produce identical outputs (reproducible for audit & compliance)",
     "\u2022 Celery Beat enables continuous monitoring with hourly automated scans across all registered targets",
     "\u2022 WebSocket real-time event pipeline: <100ms latency from scan completion to dashboard update",
     "\u2022 Docker Compose orchestration with health checks, resource limits, and automatic restart policies"],
    border=ACCENT_CYAN)

_slide_number(s, 31)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 32 — FUTURE WORK
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "FUTURE WORK", "Roadmap & Research Extensions")

future_items = [
    ("Multi-Tenant Architecture",
     "Support for multiple organizations with role-based access control, tenant isolation, and per-organization risk dashboards.",
     ACCENT_CYAN),
    ("Machine Learning Anomaly Detection",
     "Train models on baseline traffic patterns to detect zero-day threats and anomalous behavior without signature updates.",
     ACCENT_BLUE),
    ("Kubernetes Orchestration",
     "Migrate from Docker Compose to Kubernetes for horizontal scaling, auto-healing, and production-grade orchestration.",
     ACCENT_GREEN),
    ("Compliance Mapping",
     "Map vulnerabilities and risk scores to frameworks like ISO 27001, NIST CSF, and PCI-DSS for automated compliance reporting.",
     ACCENT_ORANGE),
    ("Agentic Remediation",
     "Extend the agent pipeline with a RemediationAgent that can autonomously apply fixes (firewall rules, patches, configurations) with human-in-the-loop approval.",
     ACCENT_RED),
    ("Extended Tool Integration",
     "Add support for Burp Suite API, Snyk for dependency scanning, and cloud security posture management (CSPM) tools.",
     ACCENT_CYAN),
]

for i, (title, desc, color) in enumerate(future_items):
    col = i % 2
    row = i // 2
    x = Inches(0.4) + Inches(col * 6.4)
    y = Inches(1.2) + Inches(row * 1.9)
    card = _add_shape(s, x, y, Inches(6.1), Inches(1.6),
                      fill_color=BG_CARD, border_color=color, border_width=Pt(1.5))
    tf = _tf(card)
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_font(run, size=Pt(14), color=color, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    run2 = p2.add_run()
    run2.text = desc
    _set_font(run2, size=Pt(11), color=TEXT_LIGHT)

_slide_number(s, 32)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 33 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
_add_title(s, "CONCLUSION", "Summary & Key Takeaways")

conclusions = [
    "Orchestration Security Center demonstrates that enterprise-grade security orchestration can be made accessible to SMEs through "
    "intelligent automation and AI-driven advisory generation.",

    "The Unified Risk Engine provides deterministic, reproducible risk quantification that bridges the gap between "
    "raw vulnerability data and actionable business intelligence.",

    "The multi-agent pipeline reduces mean-time-to-detection from days to minutes by autonomously chaining "
    "reconnaissance, attack, validation, and reporting without human intervention.",

    "The hybrid architecture\u2014deterministic scoring with non-deterministic AI enhancement\u2014ensures both "
    "reliability and contextual intelligence, where AI enhances but never overrides security decisions.",

    "The living lab validates the platform against a realistic SME network topology with 4 subnets and 11 "
    "vulnerable containers, proving practical effectiveness beyond theoretical design.",
]

for i, text in enumerate(conclusions):
    y = Inches(1.2) + Inches(i * 1.2)
    # Checkmark
    check = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.15), Inches(0.4), Inches(0.4))
    check.fill.solid()
    check.fill.fore_color.rgb = ACCENT_GREEN
    check.line.fill.background()
    tf = _tf(check)
    tf.margin_top = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\u2713"
    _set_font(run, size=Pt(16), color=BG_DARK, bold=True)

    _add_text_box(s, Inches(1.2), y, Inches(11.3), Inches(1.0),
                  text, size=Pt(13), color=TEXT_LIGHT)

_slide_number(s, 33)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 34 — Q&A
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()

_add_text_box(s, Inches(1), Inches(1.2), Inches(11.3), Inches(1),
              "THANK YOU", size=Pt(50), color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

_add_text_box(s, Inches(1), Inches(2.5), Inches(11.3), Inches(0.8),
              "Questions & Discussion",
              size=Pt(28), color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

# Accent line
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(3.5), Inches(3.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_CYAN
line.line.fill.background()

_add_text_box(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
              "Orchestration Security Center  \u2014  Autonomous Security Orchestration Platform for SMEs",
              size=Pt(16), color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Key topics for discussion
discussion = [
    "Deterministic vs. Non-Deterministic Architecture Decisions",
    "Unified Risk Engine Weighting Calibration",
    "Multi-Agent Orchestration Scalability",
    "Living Lab Fidelity & Real-World Applicability",
]

for i, topic in enumerate(discussion):
    y = Inches(5.0) + Inches(i * 0.45)
    _add_text_box(s, Inches(3.5), y, Inches(7), Inches(0.4),
                  f"\u25B8  {topic}", size=Pt(13), color=TEXT_LIGHT, alignment=PP_ALIGN.LEFT)

_slide_number(s, 34)


# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Found_404_Defense_Presentation.pptx")
prs.save(output_path)
print(f"\n{'='*60}")
print(f"  PRESENTATION GENERATED SUCCESSFULLY")
print(f"  Output: {output_path}")
print(f"  Slides: {len(prs.slides)}")
print(f"  Theme:  Futuristic Glass / iOS Dark")
print(f"  Font:   Merriweather Light")
print(f"{'='*60}")
