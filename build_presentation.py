"""
Build the Orchestration Security Center FYP presentation (20 slides).
Output: FINAL_PRESENTATION.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Theme ─────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0A, 0x0A, 0x14)
BG_PANEL  = RGBColor(0x12, 0x15, 0x22)
BG_PANEL2 = RGBColor(0x1A, 0x1F, 0x30)
CYAN      = RGBColor(0x00, 0xE0, 0xFF)
GREEN     = RGBColor(0x00, 0xFF, 0x88)
ORANGE    = RGBColor(0xFF, 0x9C, 0x2A)
RED       = RGBColor(0xFF, 0x4A, 0x4A)
MAGENTA   = RGBColor(0xFF, 0x4A, 0xC8)
WHITE     = RGBColor(0xEC, 0xEF, 0xF4)
GREY      = RGBColor(0x9A, 0xA3, 0xB2)
LINE      = RGBColor(0x22, 0x2A, 0x3D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────
def add_bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg

def add_rect(slide, x, y, w, h, fill=BG_PANEL, line=LINE, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.06
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, bullets, *, size=16, color=WHITE,
                bullet_color=CYAN, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet marker
        r1 = p.add_run(); r1.text = "▸  "
        r1.font.name = "Consolas"; r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = bullet_color
        # body
        r2 = p.add_run(); r2.text = b
        r2.font.name = "Calibri"; r2.font.size = Pt(size); r2.font.color.rgb = color
    return tb

def add_accent_bar(slide, x, y, h, color=CYAN, w=Inches(0.08)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.shadow.inherit = False
    return bar

def add_header(slide, title, subtitle=None, accent=CYAN):
    add_bg(slide)
    # top accent line
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.08))
    top.line.fill.background()
    top.fill.solid(); top.fill.fore_color.rgb = accent
    top.shadow.inherit = False
    # title block
    add_accent_bar(slide, Inches(0.55), Inches(0.55), Inches(0.55), accent)
    add_text(slide, Inches(0.75), Inches(0.4), Inches(11.5), Inches(0.6),
             title, size=30, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.78), Inches(0.95), Inches(11.5), Inches(0.4),
                 subtitle, size=14, color=GREY)
    # divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.4),
                                 Inches(12.2), Pt(1))
    div.line.fill.background()
    div.fill.solid(); div.fill.fore_color.rgb = LINE
    div.shadow.inherit = False
    # footer
    add_text(slide, Inches(0.55), Inches(7.05), Inches(7.5), Inches(0.35),
             "Orchestration Security Center  ·  HITU FYP 2025/2026",
             size=10, color=GREY)
    add_text(slide, Inches(10.5), Inches(7.05), Inches(2.3), Inches(0.35),
             "", size=10, color=GREY, align=PP_ALIGN.RIGHT)

def page_number(slide, n, total=20):
    add_text(slide, Inches(11.7), Inches(7.05), Inches(1.1), Inches(0.35),
             f"{n:02d} / {total:02d}",
             size=10, color=CYAN, align=PP_ALIGN.RIGHT)

def set_notes(slide, text):
    nf = slide.notes_slide.notes_text_frame
    nf.text = text

# ── Slide 1 — Title ───────────────────────────────────────────────────
def slide_01():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    # diagonal accent
    for i, col in enumerate([CYAN, MAGENTA, GREEN]):
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(-1), Inches(7.0 - i*0.12), SW + Inches(2), Pt(2))
        bar.line.fill.background()
        bar.fill.solid(); bar.fill.fore_color.rgb = col
        bar.shadow.inherit = False
    # eyebrow
    add_text(s, Inches(0.8), Inches(1.6), Inches(11), Inches(0.4),
             "FINAL YEAR PROJECT  ·  CYBERSECURITY",
             size=14, bold=True, color=CYAN)
    # title
    add_text(s, Inches(0.8), Inches(2.1), Inches(11.7), Inches(1.6),
             "Orchestration Security Center",
             size=58, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.9),
             "An Automated Security Orchestration Platform for SMEs",
             size=26, color=CYAN)
    # tagline panel
    add_rect(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.9),
             fill=BG_PANEL, line=CYAN)
    add_text(s, Inches(1.1), Inches(4.6), Inches(11.2), Inches(0.7),
             "“Turn 200 raw vulnerability alerts into 5 prioritised business actions.”",
             size=18, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # footer info
    add_text(s, Inches(0.8), Inches(6.1), Inches(6), Inches(0.4),
             "Helwan International Technological University (HITU)",
             size=14, color=GREY)
    add_text(s, Inches(0.8), Inches(6.45), Inches(6), Inches(0.4),
             "Team of 12  ·  4 sub-teams  ·  Foundational engineering focus",
             size=14, color=GREY)
    add_text(s, Inches(8.5), Inches(6.1), Inches(4), Inches(0.4),
             "Presented: July 2026",
             size=14, color=GREY, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(8.5), Inches(6.45), Inches(4), Inches(0.4),
             "Supervisor: <name>",
             size=14, color=GREY, align=PP_ALIGN.RIGHT)
    page_number(s, 1)
    set_notes(s,
        "Omar Kapil opens. Warm welcome. State the problem in one sentence: "
        "'90% of Egyptian businesses are SMEs and almost none can afford a security analyst — "
        "that is the gap we close.' Do NOT introduce the team yet; that is slide 19.")

# ── Slide 2 — The Problem ─────────────────────────────────────────────
def slide_02():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "The SME Security Gap", "Why existing approaches leave small businesses exposed")
    # left column: numbers
    add_rect(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(5.0),
             fill=BG_PANEL, line=LINE)
    add_text(s, Inches(0.85), Inches(1.85), Inches(5.4), Inches(0.5),
             "WHAT THE DATA SAYS", size=12, bold=True, color=CYAN)
    big_facts = [
        ("1 in 3", "SMEs hit by a cyberattack last year", RED),
        (">6 mo", "Avg breach cost > 6 months of revenue", ORANGE),
        ("$200K+", "Annual cost of a 3-analyst SOC", ORANGE),
        ("200+",   "Raw findings from a typical scanner run", CYAN),
    ]
    for i, (big, label, col) in enumerate(big_facts):
        y = Inches(2.4 + i*1.0)
        add_text(s, Inches(0.85), y, Inches(2.0), Inches(0.85),
                 big, size=34, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.95), y, Inches(3.4), Inches(0.85),
                 label, size=14, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # right column: what they have / what attackers see
    add_rect(s, Inches(6.85), Inches(1.7), Inches(5.95), Inches(2.4),
             fill=BG_PANEL, line=GREEN)
    add_text(s, Inches(7.1), Inches(1.85), Inches(5.5), Inches(0.4),
             "WHAT AN SME HAS", size=12, bold=True, color=GREEN)
    add_bullets(s, Inches(7.1), Inches(2.3), Inches(5.5), Inches(1.8), [
        "1 generalist IT person",
        "No SIEM, no SOC, no incident response",
        "Off-the-shelf scanners they cannot interpret",
    ], size=14, bullet_color=GREEN)

    add_rect(s, Inches(6.85), Inches(4.25), Inches(5.95), Inches(2.45),
             fill=BG_PANEL, line=RED)
    add_text(s, Inches(7.1), Inches(4.4), Inches(5.5), Inches(0.4),
             "WHAT ATTACKERS EXPLOIT", size=12, bold=True, color=RED)
    add_bullets(s, Inches(7.1), Inches(4.85), Inches(5.5), Inches(1.8), [
        "Default credentials (admin/admin123)",
        "Open management ports (RDP, SMB, Telnet)",
        "Unpatched CVEs and outdated services",
        "No log monitoring → silent compromise",
    ], size=14, bullet_color=RED)
    page_number(s, 2)
    set_notes(s,
        "Omar Kapil. Frame as cost + expertise + volume problem. Real SME breaches happen "
        "because the simple stuff is unmonitored. End with: 'Even if they buy a scanner, "
        "they cannot read the output.' Bridge to slide 3.")

# ── Slide 3 — Existing Solutions Don't Fit ────────────────────────────
def slide_03():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Why Off-the-Shelf Tools Don’t Fit",
               "We respect them — we just need a different layer on top")
    rows = [
        ("Tenable / Rapid7 / Qualys", "Strong scanning",
         "Expensive · cloud-only · no remediation guidance", RED),
        ("OpenVAS / Nmap (raw)", "Free and powerful",
         "Output requires expert interpretation", ORANGE),
        ("Pure-LLM ‘AI security’ startups", "Modern UX",
         "Hallucinate · non-deterministic · unauditable", RED),
        ("OUR ANSWER: Orchestration Security Center", "Free + explainable + deterministic",
         "Hybrid: math for scoring, LLM only for English", GREEN),
    ]
    y = Inches(1.75)
    for name, pro, con, col in rows:
        add_rect(s, Inches(0.55), y, Inches(12.2), Inches(1.15),
                 fill=BG_PANEL, line=col)
        add_accent_bar(s, Inches(0.55), y, Inches(1.15), col)
        add_text(s, Inches(0.8), y, Inches(4.3), Inches(1.15),
                 name, size=15, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.2), y, Inches(3.5), Inches(1.15),
                 "✓  " + pro, size=13, color=GREEN,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.8), y, Inches(3.9), Inches(1.15),
                 "✗  " + con if col != GREEN else "→  " + con,
                 size=13, color=col, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.27)
    page_number(s, 3)
    set_notes(s,
        "Omar Kapil. Acknowledge the existing tools are good — we use Nmap and Nuclei "
        "underneath. Innovation is the orchestration layer. Set up slide 4.")

# ── Slide 4 — Our Solution ────────────────────────────────────────────
def slide_04():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Our Solution",
               "The orchestration layer for SME security")
    # left
    add_rect(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(5.0),
             fill=BG_PANEL, line=CYAN)
    add_text(s, Inches(0.85), Inches(1.85), Inches(5.5), Inches(0.4),
             "WHAT IT DOES", size=12, bold=True, color=CYAN)
    add_bullets(s, Inches(0.85), Inches(2.3), Inches(5.5), Inches(4.3), [
        "Automated DAST platform — discover · attack · validate · score · report",
        "4-stage AI agent pipeline runs unattended",
        "Deterministic 0–100 risk score (no LLM in scoring path)",
        "Gemini LLM advisory only — plain-English explanations",
        "Single-laptop deploy: docker compose up",
        "5 prioritised actions, not 200 raw alerts",
    ], size=15, bullet_color=CYAN)

    # right
    add_rect(s, Inches(6.85), Inches(1.7), Inches(5.95), Inches(5.0),
             fill=BG_PANEL2, line=LINE)
    add_text(s, Inches(7.1), Inches(1.85), Inches(5.5), Inches(0.4),
             "BEFORE  →  AFTER", size=12, bold=True, color=GREEN)
    # before card
    add_rect(s, Inches(7.1), Inches(2.3), Inches(5.45), Inches(2.0),
             fill=BG_DARK, line=RED)
    add_text(s, Inches(7.3), Inches(2.4), Inches(5.0), Inches(0.4),
             "BEFORE", size=11, bold=True, color=RED)
    add_text(s, Inches(7.3), Inches(2.75), Inches(5.0), Inches(1.5),
             "200 CSV rows of CVEs.\nNo priority. No business context.\nNobody reads it.",
             size=14, color=WHITE)
    # after card
    add_rect(s, Inches(7.1), Inches(4.55), Inches(5.45), Inches(2.05),
             fill=BG_DARK, line=GREEN)
    add_text(s, Inches(7.3), Inches(4.65), Inches(5.0), Inches(0.4),
             "AFTER", size=11, bold=True, color=GREEN)
    add_text(s, Inches(7.3), Inches(5.0), Inches(5.0), Inches(1.5),
             "5 ranked action cards.\nPlain-English impact + remediation.\nThe owner knows what to fix first.",
             size=14, color=WHITE)
    page_number(s, 4)
    set_notes(s,
        "Omar Kapil. Thesis statement slide. 'We are not a scanner. We are the brain "
        "that orchestrates scanners and translates output into action.' Hand off to Reem.")

# ── Slide 5 — 4 Pillars ───────────────────────────────────────────────
def slide_05():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Core Innovation — Four Pillars",
               "What makes Orchestration Security Center different")
    pillars = [
        ("01", "Deterministic Tool Chaining", CYAN,
         "Nmap result drives Nuclei template choice.\nPort 445 → SMB templates only — never SQLi.\nFewer false positives, faster scans.",
         "nuclei_wrapper.py"),
        ("02", "Hybrid Intelligence", MAGENTA,
         "Math-based scoring + AI-based explanation.\nDeterministic decisions; LLM only for English.\nAuditable + accessible.",
         "unified_risk_engine.py"),
        ("03", "Tamper-Evident Audit", GREEN,
         "Every agent action hash-chained.\nsha256(prev_hash + entry).\nReplayable like a blockchain.",
         "agent_orchestrator.py · BaseAgent.log_action()"),
        ("04", "Single-Pane Visibility", ORANGE,
         "Scanner findings + Wazuh SIEM alerts fused.\nOne dashboard, one source of truth.\nReact + WebSocket-driven UI.",
         "alert_correlator.py · RealTimeContext.jsx"),
    ]
    positions = [
        (Inches(0.55), Inches(1.7)),
        (Inches(6.95), Inches(1.7)),
        (Inches(0.55), Inches(4.45)),
        (Inches(6.95), Inches(4.45)),
    ]
    for (num, title, col, body, file_), (x, y) in zip(pillars, positions):
        add_rect(s, x, y, Inches(5.85), Inches(2.6), fill=BG_PANEL, line=col)
        add_accent_bar(s, x, y, Inches(2.6), col, w=Inches(0.1))
        add_text(s, x + Inches(0.3), y + Inches(0.15), Inches(1.0), Inches(0.5),
                 num, size=22, bold=True, color=col)
        add_text(s, x + Inches(1.1), y + Inches(0.15), Inches(4.6), Inches(0.5),
                 title, size=18, bold=True, color=WHITE)
        add_text(s, x + Inches(0.3), y + Inches(0.85), Inches(5.4), Inches(1.3),
                 body, size=12, color=WHITE)
        add_text(s, x + Inches(0.3), y + Inches(2.15), Inches(5.4), Inches(0.4),
                 "▣  " + file_, size=10, color=col, font="Consolas")
    page_number(s, 5)
    set_notes(s,
        "Reem Amin takes over here. Emphasise determinism — committee always asks "
        "'how do you trust the AI?'. Answer: we do not trust it for decisions; only "
        "for English. The math is auditable.")

# ── Slide 6 — Architecture ────────────────────────────────────────────
def slide_06():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "System Architecture",
               "Single-laptop deployment · 7 services · 4-zone lab")
    # User
    def box(x, y, w, h, label, sub, col):
        add_rect(s, x, y, w, h, fill=BG_PANEL, line=col)
        add_text(s, x, y + Inches(0.15), w, Inches(0.4),
                 label, size=13, bold=True, color=col,
                 align=PP_ALIGN.CENTER)
        if sub:
            add_text(s, x, y + Inches(0.55), w, Inches(0.35),
                     sub, size=10, color=GREY, align=PP_ALIGN.CENTER, font="Consolas")
        return (x, y, w, h)

    def arrow(x1, y1, x2, y2, col=CYAN, weight=1.25):
        ln = s.shapes.add_connector(2, x1, y1, x2, y2)  # straight
        ln.line.color.rgb = col
        ln.line.width = Pt(weight)

    # Row 1 — User & Caddy
    b_user = box(Inches(0.7), Inches(1.7), Inches(2.2), Inches(0.95),
                 "User Browser", "https://localhost", WHITE)
    b_caddy = box(Inches(3.4), Inches(1.7), Inches(2.4), Inches(0.95),
                  "Caddy TLS Proxy", "auto-TLS · :443", CYAN)
    # Row 2 — Frontend / Backend
    b_fe = box(Inches(0.7), Inches(3.0), Inches(2.6), Inches(1.0),
               "React Frontend", "Vite · :5173", MAGENTA)
    b_be = box(Inches(3.6), Inches(3.0), Inches(2.6), Inches(1.0),
               "FastAPI Backend", "Python 3.11 · :8000", CYAN)
    # Row 3 — Data
    b_db = box(Inches(0.7), Inches(4.45), Inches(2.0), Inches(0.9),
               "PostgreSQL", "ACID state · :5432", GREEN)
    b_redis = box(Inches(2.85), Inches(4.45), Inches(2.0), Inches(0.9),
                  "Redis", "queue + pub/sub", ORANGE)
    b_celery = box(Inches(5.0), Inches(4.45), Inches(2.0), Inches(0.9),
                   "Celery Worker", "async scans", ORANGE)
    # Row 4 — Scanners + AI
    b_nmap = box(Inches(0.7), Inches(5.7), Inches(1.55), Inches(0.85),
                 "Nmap", "ports · OS · NSE", WHITE)
    b_nuc  = box(Inches(2.35), Inches(5.7), Inches(1.55), Inches(0.85),
                 "Nuclei", "YAML CVE templates", WHITE)
    b_ovas = box(Inches(4.0), Inches(5.7), Inches(1.55), Inches(0.85),
                 "OpenVAS", "deep CVE scan", WHITE)
    b_gem  = box(Inches(5.65), Inches(5.7), Inches(1.55), Inches(0.85),
                 "Gemini LLM", "advisory only", MAGENTA)

    # Right column — SIEM + Lab
    add_rect(s, Inches(7.45), Inches(1.7), Inches(5.3), Inches(2.5),
             fill=BG_PANEL2, line=GREEN)
    add_text(s, Inches(7.55), Inches(1.8), Inches(5.0), Inches(0.4),
             "SIEM TIER", size=12, bold=True, color=GREEN)
    add_bullets(s, Inches(7.6), Inches(2.2), Inches(5.0), Inches(2.0), [
        "Wazuh (HIDS + rules)",
        "Elasticsearch (log index)",
        "Kibana (dashboards)",
        "alert_correlator.py → DB",
    ], size=13, bullet_color=GREEN)

    add_rect(s, Inches(7.45), Inches(4.35), Inches(5.3), Inches(2.4),
             fill=BG_PANEL2, line=RED)
    add_text(s, Inches(7.55), Inches(4.45), Inches(5.0), Inches(0.4),
             "LAB ENVIRONMENT (internal: true)", size=12, bold=True, color=RED)
    add_bullets(s, Inches(7.6), Inches(4.85), Inches(5.0), Inches(2.0), [
        "DMZ — Juice Shop · API GW · DNS",
        "CORP — Samba · GreenMail · workstation",
        "DATA — PostgreSQL · Redis (no auth)",
        "MGMT — traffic gen · log shipper",
    ], size=13, bullet_color=RED)

    # connectors (key flow)
    arrow(Inches(1.8), Inches(2.65), Inches(4.6), Inches(2.65), CYAN)  # user -> caddy
    arrow(Inches(4.6), Inches(2.65), Inches(2.0), Inches(3.0), CYAN)
    arrow(Inches(4.6), Inches(2.65), Inches(4.9), Inches(3.0), CYAN)
    arrow(Inches(4.9), Inches(4.0), Inches(1.7), Inches(4.45), CYAN)
    arrow(Inches(4.9), Inches(4.0), Inches(3.85), Inches(4.45), ORANGE)
    arrow(Inches(3.85), Inches(5.35), Inches(6.0), Inches(5.35), ORANGE)
    arrow(Inches(6.0), Inches(5.35), Inches(6.0), Inches(5.7), ORANGE)
    page_number(s, 6)
    set_notes(s,
        "Reem Amin. Walk left-to-right along the data flow. Pause on Redis pub/sub: "
        "'this is how a 3-minute scan streams live to the browser without polling'. "
        "Mention the whole stack runs on 8 GB RAM.")

# ── Slide 7 — Tech Stack ──────────────────────────────────────────────
def slide_07():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Tech Stack — Why Each Choice",
               "Right tool for each layer · not one tool for everything")
    headers = ["LAYER", "TECHNOLOGY", "WHY THIS CHOICE"]
    rows = [
        ("Frontend",   "React 18 · Vite · Tailwind · D3.js",   "Lazy-loaded panels · D3 force graph · cyber theme"),
        ("State",      "React Query · useReducer (RealTime)",  "Server cache + WebSocket reducer pattern"),
        ("Backend",    "FastAPI · Pydantic v2 · SQLAlchemy 2", "Auto-OpenAPI · async-first · type-safe"),
        ("Async",      "Celery + Redis (broker + pub/sub)",    "Fan-out scans · live event bridge"),
        ("Storage",    "PostgreSQL + Alembic",                 "ACID for security-critical state"),
        ("Scanners",   "Nmap · Nuclei · OpenVAS",              "Discovery → templates → deep CVE"),
        ("SIEM",       "Wazuh · Elasticsearch · Kibana",       "Open-source, production-grade"),
        ("AI",         "Google Gemini 2.0 Flash",              "Cheap · fast · JSON-mode advisory"),
        ("DevOps",     "Docker Compose · Caddy · GH Actions · Trivy",
                                                                 "Reproducible · TLS · scanned in CI"),
    ]
    # header row
    y = Inches(1.7)
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.5),
             fill=BG_PANEL2, line=CYAN)
    add_text(s, Inches(0.75), y, Inches(2.4), Inches(0.5),
             headers[0], size=12, bold=True, color=CYAN, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.2), y, Inches(4.2), Inches(0.5),
             headers[1], size=12, bold=True, color=CYAN, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.5), y, Inches(5.1), Inches(0.5),
             headers[2], size=12, bold=True, color=CYAN, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.55)
    for layer, tech, why in rows:
        add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.48),
                 fill=BG_PANEL, line=LINE)
        add_text(s, Inches(0.75), y, Inches(2.4), Inches(0.48),
                 layer, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.2), y, Inches(4.2), Inches(0.48),
                 tech, size=11, color=CYAN, anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
        add_text(s, Inches(7.5), y, Inches(5.1), Inches(0.48),
                 why, size=11, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.51)
    page_number(s, 7)
    set_notes(s,
        "Reem. Don't read the table. Emphasise the philosophy: Python/FastAPI for "
        "orchestration logic (the brain), C/Go scanners (Nmap/Nuclei) for low-level "
        "work (the muscle), React for visualisation.")

# ── Slide 8 — Frontend ────────────────────────────────────────────────
def slide_08():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Module Deep Dive — Frontend",
               "React 18 · WebSocket-first · auth-aware UI", accent=MAGENTA)
    add_rect(s, Inches(0.55), Inches(1.7), Inches(6.5), Inches(5.0),
             fill=BG_PANEL, line=MAGENTA)
    add_text(s, Inches(0.85), Inches(1.85), Inches(6.0), Inches(0.4),
             "ARCHITECTURE & TECH", size=12, bold=True, color=MAGENTA)
    add_bullets(s, Inches(0.85), Inches(2.3), Inches(6.0), Inches(4.4), [
        "Vite-bundled · code-split · 25+ dashboard panels",
        "RealTimeContext.jsx — useReducer + WebSocket envelope {type, payload, seq, ts}",
        "Exponential-backoff auto-reconnect (1s → 30s cap)",
        "react-force-graph-2d — node colour by risk score",
        "D3.js heatmap · Chart.js trends · Recharts gauges",
        "JWT in localStorage · ProtectedRoute · RoleGuard",
        "TailwindCSS cyber theme (cyber-cyan / cyber-red / glass cards)",
    ], size=14, bullet_color=MAGENTA)

    # right: feature highlights
    add_rect(s, Inches(7.35), Inches(1.7), Inches(5.4), Inches(2.4),
             fill=BG_PANEL2, line=CYAN)
    add_text(s, Inches(7.55), Inches(1.85), Inches(5.0), Inches(0.4),
             "REAL-TIME EVENT LOOP", size=12, bold=True, color=CYAN)
    add_bullets(s, Inches(7.55), Inches(2.3), Inches(5.0), Inches(1.8), [
        "WS open → SET_CONNECTED",
        "RISK_UPDATE → KPI cards re-render",
        "LOG_STREAM → OrchestrationFeed (max 200)",
        "ALERT_NEW → toast + alert badge",
    ], size=12, bullet_color=CYAN)

    add_rect(s, Inches(7.35), Inches(4.25), Inches(5.4), Inches(2.45),
             fill=BG_PANEL2, line=GREEN)
    add_text(s, Inches(7.55), Inches(4.4), Inches(5.0), Inches(0.4),
             "AUTH-AWARE COMPONENTS", size=12, bold=True, color=GREEN)
    add_bullets(s, Inches(7.55), Inches(4.85), Inches(5.0), Inches(1.8), [
        "AuthContext stores JWT + role",
        "<ProtectedRoute /> redirects guests",
        "<RoleGuard role='admin' /> hides UI",
        "Force-password-change on first login",
    ], size=12, bullet_color=GREEN)
    page_number(s, 8)
    set_notes(s,
        "Marize Ehap takes the mic. Switch to live dashboard; click a topology node; "
        "narrate that the hover tooltip is fed from the same Redis-backed RealTimeContext "
        "as the KPI cards — single source of truth.")

# ── Slide 9 — Backend & Auth ──────────────────────────────────────────
def slide_09():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Module Deep Dive — Backend API & Auth",
               "FastAPI · JWT · RBAC · audit chain")
    # Left
    add_rect(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(5.0),
             fill=BG_PANEL, line=CYAN)
    add_text(s, Inches(0.85), Inches(1.85), Inches(5.5), Inches(0.4),
             "FACTS", size=12, bold=True, color=CYAN)
    add_bullets(s, Inches(0.85), Inches(2.3), Inches(5.5), Inches(4.3), [
        "40+ REST endpoints · auto-OpenAPI at /docs",
        "13 router groups (auth · scans · targets · findings · …)",
        "bcrypt password hash · JWT (HS256) · Depends(get_current_user)",
        "RBAC: admin · analyst · viewer · require_role()",
        "Fernet-encrypted credentials at rest",
        "Lifespan: Alembic check · scan-reaper · admin-seed · Redis listener",
        "Hash-chained AgentLog → tamper-evident",
    ], size=14, bullet_color=CYAN)

    # Right: sequence diagram
    add_rect(s, Inches(6.85), Inches(1.7), Inches(5.95), Inches(5.0),
             fill=BG_PANEL2, line=CYAN)
    add_text(s, Inches(7.05), Inches(1.85), Inches(5.5), Inches(0.4),
             "AUTH SEQUENCE", size=12, bold=True, color=CYAN)
    seq = [
        ("Browser", "POST /auth/login", "Email + password"),
        ("FastAPI",  "verify_password()", "bcrypt compare"),
        ("FastAPI",  "create_access_token()", "JWT HS256"),
        ("Browser",  "Authorization: Bearer …", "Stored in localStorage"),
        ("FastAPI",  "Depends(get_current_user)", "Decode + DB lookup"),
        ("Endpoint", "require_role('admin')",   "RBAC guard"),
    ]
    y = Inches(2.3)
    for actor, call, note in seq:
        add_rect(s, Inches(7.05), y, Inches(1.4), Inches(0.55),
                 fill=BG_DARK, line=CYAN)
        add_text(s, Inches(7.05), y, Inches(1.4), Inches(0.55),
                 actor, size=11, bold=True, color=CYAN,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(8.55), y, Inches(2.6), Inches(0.55),
                 call, size=11, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
                 font="Consolas")
        add_text(s, Inches(11.2), y, Inches(1.5), Inches(0.55),
                 note, size=10, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.7)
    page_number(s, 9)
    set_notes(s,
        "Reem Amin. Show /docs live, expand /scans/ai. Explain the Depends() pattern: "
        "this is how 40 endpoints share one auth check without copy-paste.")

# ── Slide 10 — AI Agent Pipeline ──────────────────────────────────────
def slide_10():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Module Deep Dive — The AI Agent Pipeline",
               "agent_orchestrator.py · 4 stages · async + rate-limited", accent=MAGENTA)
    # Pipeline boxes across the top
    stages = [
        ("RECON",     "Nmap · Playwright crawl",   CYAN),
        ("ATTACK",    "httpx payloads (SQLi/XSS/BOLA/SSRF)", ORANGE),
        ("VALIDATION","drop confidence < 0.6",     GREEN),
        ("REPORTING", "markdown + PDF export",     MAGENTA),
    ]
    x = Inches(0.55); y = Inches(1.75); w = Inches(2.85); h = Inches(1.6)
    for i, (name, sub, col) in enumerate(stages):
        bx = x + Inches(i * 3.06)
        add_rect(s, bx, y, w, h, fill=BG_PANEL, line=col)
        add_accent_bar(s, bx, y, h, col, w=Inches(0.08))
        add_text(s, bx + Inches(0.2), y + Inches(0.15), w, Inches(0.55),
                 f"STAGE {i+1}", size=11, bold=True, color=col)
        add_text(s, bx + Inches(0.2), y + Inches(0.55), w, Inches(0.5),
                 name, size=20, bold=True, color=WHITE)
        add_text(s, bx + Inches(0.2), y + Inches(1.05), w - Inches(0.3), Inches(0.55),
                 sub, size=11, color=GREY)
        if i < 3:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                bx + w + Inches(0.02), y + Inches(0.65),
                Inches(0.18), Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = CYAN
            arr.line.fill.background()
            arr.shadow.inherit = False

    # Lower split: BaseAgent + Advisory side-channel
    add_rect(s, Inches(0.55), Inches(3.6), Inches(6.0), Inches(3.1),
             fill=BG_PANEL2, line=CYAN)
    add_text(s, Inches(0.8), Inches(3.75), Inches(5.5), Inches(0.4),
             "BaseAgent (abstract)", size=12, bold=True, color=CYAN, font="Consolas")
    add_bullets(s, Inches(0.8), Inches(4.2), Inches(5.5), Inches(2.4), [
        "AsyncLimiter — 10 rps per agent",
        "log_action() → AgentLog · sha256 hash chain",
        "llm_reason() → Gemini 2.0 Flash (fallback if no key)",
        "Async DB session per Celery task (own event loop)",
        "execute(context) → enriched context to next stage",
    ], size=13, bullet_color=CYAN)

    add_rect(s, Inches(6.85), Inches(3.6), Inches(5.95), Inches(3.1),
             fill=BG_PANEL2, line=MAGENTA)
    add_text(s, Inches(7.1), Inches(3.75), Inches(5.5), Inches(0.4),
             "IntelligenceAgent — Advisory Side-Channel", size=12, bold=True,
             color=MAGENTA)
    add_bullets(s, Inches(7.1), Inches(4.2), Inches(5.5), Inches(2.4), [
        "NEVER on the decision path",
        "Receives finished scan results",
        "Asks Gemini for plain-English impact + remediation",
        "JSON-mode → structured output for the dashboard",
        "llm_guard.py blocks destructive output",
    ], size=13, bullet_color=MAGENTA)
    page_number(s, 10)
    set_notes(s,
        "Yousef Abdel Hady supports Reem here. Each agent is async, rate-limited, "
        "and writes a tamper-proof log. Replay any scan from AgentLog. Open the "
        "Orchestration Feed during demo so the committee sees the pipeline live.")

# ── Slide 11 — Risk Engine ────────────────────────────────────────────
def slide_11():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Module Deep Dive — Risk Engine (Math, Not Magic)",
               "UnifiedRiskEngine · explainable 0–100 score")
    # Formula bar
    add_rect(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(0.85),
             fill=BG_PANEL2, line=CYAN)
    add_text(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(0.85),
             "score = ((Σ severity·confidence) + Σ port_penalties) × asset_mult × exposure_mod   →   capped at 100",
             size=15, bold=True, color=CYAN, font="Consolas",
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # 3 columns of weights
    cols = [
        ("SEVERITY WEIGHTS", RED, [
            "CRITICAL  +25 × confidence",
            "HIGH         +15 × confidence",
            "MEDIUM    +7 × confidence",
            "LOW           +2 × confidence",
            "INFO          0",
        ]),
        ("PORT PENALTIES", ORANGE, [
            "Telnet 23      +20",
            "SMB 445       +20",
            "FTP 21          +15",
            "RDP 3389     +15",
            "Redis 6379  +10  ·  Postgres 5432  +10",
        ]),
        ("MODIFIERS", GREEN, [
            "Asset CRITICAL  × 1.5",
            "Asset LOW             × 0.8",
            "Internal IP             × 0.6",
            "Public IP                × 1.0",
            "Health = 100 − Risk",
        ]),
    ]
    x = Inches(0.55); y = Inches(2.7)
    for i, (h, col, items) in enumerate(cols):
        bx = x + Inches(i * 4.1)
        add_rect(s, bx, y, Inches(4.0), Inches(2.6), fill=BG_PANEL, line=col)
        add_text(s, bx + Inches(0.2), y + Inches(0.1), Inches(3.6), Inches(0.4),
                 h, size=12, bold=True, color=col)
        add_bullets(s, bx + Inches(0.2), y + Inches(0.55), Inches(3.6), Inches(2.0),
                    items, size=11, bullet_color=col)
    # Worked example
    y2 = Inches(5.45)
    add_rect(s, Inches(0.55), y2, Inches(12.2), Inches(1.3),
             fill=BG_PANEL2, line=GREEN)
    add_text(s, Inches(0.75), y2 + Inches(0.1), Inches(12.0), Inches(0.4),
             "WORKED EXAMPLE — Juice Shop on public IP", size=12, bold=True, color=GREEN)
    add_text(s, Inches(0.75), y2 + Inches(0.5), Inches(12.0), Inches(0.7),
             "2 × CRITICAL (50) + 5 × HIGH (75) + port 80 open + public IP × 1.0 → "
             "score capped at 100 → reported as 95/100  ·  no LLM involved",
             size=13, color=WHITE, font="Consolas")
    page_number(s, 11)
    set_notes(s,
        "Yousef. Committee always asks: 'is your AI a black box?'. Answer: 'No — "
        "the score is pure arithmetic. The LLM only writes the English summary.'")

# ── Slide 12 — Scanning Engine ────────────────────────────────────────
def slide_12():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Module Deep Dive — Scanning Engine",
               "Deterministic tool chaining · Nmap → Nuclei → OpenVAS", accent=ORANGE)
    # Decision tree visual
    add_rect(s, Inches(0.55), Inches(1.7), Inches(7.5), Inches(5.0),
             fill=BG_PANEL, line=ORANGE)
    add_text(s, Inches(0.8), Inches(1.85), Inches(7.0), Inches(0.4),
             "DETERMINISTIC CHAINING", size=12, bold=True, color=ORANGE)
    # central nmap node
    add_rect(s, Inches(0.95), Inches(2.4), Inches(2.1), Inches(0.7),
             fill=BG_DARK, line=ORANGE)
    add_text(s, Inches(0.95), Inches(2.4), Inches(2.1), Inches(0.7),
             "Nmap  →  ports", size=13, bold=True, color=ORANGE,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, font="Consolas")
    # Branches
    branches = [
        ("port 80 / 443",   "→ http-sqli · http-xss · bola · ssrf", CYAN),
        ("port 445 (SMB)",  "→ smb-enum · smb-default-creds",      MAGENTA),
        ("port 6379 (Redis)","→ redis-unauth · redis-rce",          GREEN),
        ("port 5432 (PG)",  "→ postgres-weak-creds",               WHITE),
        ("port 21 (FTP)",   "→ ftp-anon · ftp-weak-creds",         ORANGE),
    ]
    y = Inches(3.25)
    for cond, action, col in branches:
        add_text(s, Inches(0.85), y, Inches(2.3), Inches(0.45),
                 cond, size=12, bold=True, color=col, font="Consolas")
        add_text(s, Inches(3.2), y, Inches(4.7), Inches(0.45),
                 action, size=11, color=WHITE, font="Consolas")
        y += Inches(0.55)

    # Right column
    add_rect(s, Inches(8.4), Inches(1.7), Inches(4.4), Inches(2.4),
             fill=BG_PANEL2, line=CYAN)
    add_text(s, Inches(8.55), Inches(1.85), Inches(4.0), Inches(0.4),
             "NMAP MODES", size=12, bold=True, color=CYAN)
    add_bullets(s, Inches(8.55), Inches(2.3), Inches(4.0), Inches(1.8), [
        "quick → -sV -F -T4",
        "full   → -sV -O -T4",
        "deep  → -sV -O -A --script=vulners,smb-os-discovery,http-enum",
    ], size=11, bullet_color=CYAN)

    add_rect(s, Inches(8.4), Inches(4.25), Inches(4.4), Inches(2.45),
             fill=BG_PANEL2, line=GREEN)
    add_text(s, Inches(8.55), Inches(4.4), Inches(4.0), Inches(0.4),
             "GUARDRAILS", size=12, bold=True, color=GREEN)
    add_bullets(s, Inches(8.55), Inches(4.85), Inches(4.0), Inches(1.8), [
        "scope_guard.py — allow-list",
        "Out-of-scope rejected before any packet",
        "finding_dedup.py — collapses repeats",
        "Audit-logged for every rejection",
    ], size=11, bullet_color=GREEN)
    page_number(s, 12)
    set_notes(s,
        "Shahd Paher takes over. 'Most platforms run every template against every target — "
        "that's why they generate 200 alerts. We chain deterministically. Same coverage, "
        "fewer false positives, faster scans.'")

# ── Slide 13 — SIEM ───────────────────────────────────────────────────
def slide_13():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Module Deep Dive — SIEM & Log Analytics",
               "Wazuh + Elasticsearch + Kibana · one-pane visibility", accent=GREEN)
    # Left
    add_rect(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(5.0),
             fill=BG_PANEL, line=GREEN)
    add_text(s, Inches(0.85), Inches(1.85), Inches(5.5), Inches(0.4),
             "INTEGRATIONS", size=12, bold=True, color=GREEN)
    add_bullets(s, Inches(0.85), Inches(2.3), Inches(5.5), Inches(4.3), [
        "wazuh_integration.py — Basic→JWT · httpx async · /alerts",
        "elastic_integration.py — DSL: bool · must · range · match",
        "alert_correlator.py — Wazuh alert ↔ scan vuln (IP + port + time)",
        "Custom Wazuh rules in lab/wazuh/custom_rules.xml",
        "Pre-built Kibana dashboard sme_overview.ndjson",
    ], size=14, bullet_color=GREEN)

    # Right — pipeline
    add_rect(s, Inches(6.85), Inches(1.7), Inches(5.95), Inches(5.0),
             fill=BG_PANEL2, line=CYAN)
    add_text(s, Inches(7.05), Inches(1.85), Inches(5.5), Inches(0.4),
             "DATA FLOW (lab → dashboard)", size=12, bold=True, color=CYAN)
    flow = [
        "lab containers emit access/error logs",
        "log_shipper.py — every 10 s",
        "Elasticsearch index sme-lab-*",
        "Wazuh rule engine fires alert",
        "GET /api/v1/siem/alerts (poll 30 s)",
        "SIEM tab renders alert + correlated vuln",
    ]
    y = Inches(2.3)
    for i, step in enumerate(flow):
        add_rect(s, Inches(7.05), y, Inches(0.55), Inches(0.55),
                 fill=BG_DARK, line=CYAN)
        add_text(s, Inches(7.05), y, Inches(0.55), Inches(0.55),
                 str(i+1), size=14, bold=True, color=CYAN,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, font="Consolas")
        add_text(s, Inches(7.75), y, Inches(5.0), Inches(0.55),
                 step, size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.7)
    page_number(s, 13)
    set_notes(s,
        "Mariz Ehap (or Shahd narrates). 'A scanner finding plus a SIEM alert on the "
        "same IP equals high confidence — that is correlation.' This is the slide that "
        "answers 'what about runtime threats, not just vulnerabilities?'.")

# ── Slide 14 — Async pipeline ─────────────────────────────────────────
def slide_14():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Module Deep Dive — Async Pipeline & Real-Time Bridge",
               "Celery · Redis pub/sub · WebSocket")
    # Steps
    steps = [
        ("1", "Click Scan",   "POST /scans/ai (FastAPI)",        CYAN),
        ("2", "Enqueue task", "celery_app.send_task('run_full_scan')", CYAN),
        ("3", "Worker pops",  "fresh asyncio loop per task",     ORANGE),
        ("4", "Run scanners", "Nmap → Nuclei → Risk Engine",     ORANGE),
        ("5", "Publish event","Redis PUBLISH ws_events",         MAGENTA),
        ("6", "Bridge listens","main.py redis_event_listener()", MAGENTA),
        ("7", "Broadcast",    "WebSocketManager.send_all(env)",  GREEN),
        ("8", "React update", "useReducer dispatch UPDATE_RISK", GREEN),
    ]
    x = Inches(0.55); y = Inches(1.75); w = Inches(6.0); h = Inches(0.6)
    for i, (n, what, code, col) in enumerate(steps):
        row_y = y + Inches(i * 0.65)
        col_x = x if i < 4 else x + Inches(6.3)
        row_y = y + Inches((i % 4) * 0.65) if i >= 4 else row_y
        add_rect(s, col_x, row_y, w, h, fill=BG_PANEL, line=col)
        add_rect(s, col_x, row_y, Inches(0.6), h, fill=col, line=col)
        add_text(s, col_x, row_y, Inches(0.6), h,
                 n, size=18, bold=True, color=BG_DARK,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        add_text(s, col_x + Inches(0.75), row_y, Inches(2.1), h,
                 what, size=12, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, col_x + Inches(2.85), row_y, Inches(3.0), h,
                 code, size=10, color=GREY, font="Consolas",
                 anchor=MSO_ANCHOR.MIDDLE)

    # Footer note
    add_rect(s, Inches(0.55), Inches(5.9), Inches(12.2), Inches(0.85),
             fill=BG_PANEL2, line=CYAN)
    add_text(s, Inches(0.75), Inches(5.9), Inches(12.0), Inches(0.85),
             "scan_reaper.py — kills any scan stuck in RUNNING > 10 min on startup    "
             "·    WebSocket envelope: {type, payload, seq, ts}",
             size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER, font="Consolas")
    page_number(s, 14)
    set_notes(s,
        "Mohamed Shaban (supports) / Reem narrates. 'This is the plumbing. Without it, "
        "the dashboard would just spin. With it, you watch the AI think in real time.'")

# ── Slide 15 — DevOps ─────────────────────────────────────────────────
def slide_15():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "DevOps & Infrastructure",
               "docker compose · Caddy · GitHub Actions · Trivy", accent=GREEN)
    # Left — services
    add_rect(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(5.0),
             fill=BG_PANEL, line=GREEN)
    add_text(s, Inches(0.85), Inches(1.85), Inches(5.5), Inches(0.4),
             "SEVEN SERVICES (docker-compose.yml)", size=12, bold=True, color=GREEN)
    services = [
        ("frontend",       "React 18 · Vite · 5173",  CYAN),
        ("backend",        "FastAPI · 8000",          CYAN),
        ("celery_worker",  "scan executor",           ORANGE),
        ("celery_beat",    "scheduled jobs",          ORANGE),
        ("postgres",       "ACID state · 5432",       GREEN),
        ("redis",          "queue + pub/sub · 6379",  ORANGE),
        ("caddy",          "TLS reverse proxy · 443", MAGENTA),
    ]
    y = Inches(2.3)
    for name, sub, col in services:
        add_rect(s, Inches(0.85), y, Inches(5.5), Inches(0.45),
                 fill=BG_DARK, line=col)
        add_text(s, Inches(1.0), y, Inches(2.3), Inches(0.45),
                 name, size=12, bold=True, color=col,
                 anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
        add_text(s, Inches(3.4), y, Inches(2.9), Inches(0.45),
                 sub, size=11, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.55)

    # Right — CI / lab / startup
    add_rect(s, Inches(6.85), Inches(1.7), Inches(5.95), Inches(2.4),
             fill=BG_PANEL2, line=CYAN)
    add_text(s, Inches(7.05), Inches(1.85), Inches(5.5), Inches(0.4),
             "CI / CD", size=12, bold=True, color=CYAN)
    add_bullets(s, Inches(7.05), Inches(2.3), Inches(5.5), Inches(1.8), [
        ".github/workflows/ci.yml — pytest + Trivy on every PR",
        ".github/workflows/cd.yml — auto-deploy on main merge",
        "Trivy gate rejects HIGH/CRITICAL CVE images",
    ], size=12, bullet_color=CYAN)

    add_rect(s, Inches(6.85), Inches(4.25), Inches(5.95), Inches(2.45),
             fill=BG_PANEL2, line=ORANGE)
    add_text(s, Inches(7.05), Inches(4.4), Inches(5.5), Inches(0.4),
             "ONE-COMMAND DEPLOY", size=12, bold=True, color=ORANGE)
    add_bullets(s, Inches(7.05), Inches(4.85), Inches(5.5), Inches(1.8), [
        "start-lite.ps1 — full stack in <4 min",
        "lab_setup.ps1 start|stop|seed|reset",
        "8 GB RAM budget · runs on a laptop",
        "Caddy: self-signed (dev) · Let's Encrypt (prod)",
    ], size=12, bullet_color=ORANGE)
    page_number(s, 15)
    set_notes(s,
        "Omar Kapil retakes the mic. 'Reproducibility was a hard requirement. Any "
        "examiner can clone the repo, run start-lite.ps1, and have the platform live "
        "in 4 minutes.'")

# ── Slide 16 — Securing the Platform ──────────────────────────────────
def slide_16():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Securing the Platform Itself",
               "SECURITY_AUDIT.md · we eat our own dog food", accent=RED)
    items = [
        ("AUTH",     "bcrypt · JWT · RBAC · force-password-change", CYAN),
        ("TRANSPORT","Caddy auto-TLS · HSTS · secure cookies",      MAGENTA),
        ("INPUT",    "Pydantic strict · scope_guard allow-list · no shell concat", GREEN),
        ("AI",       "llm_guard blocks destructive output · token budget · LLM-disabled mode tested", ORANGE),
        ("SUPPLY",   "Trivy gate in CI rejects HIGH/CRITICAL CVE images", RED),
        ("AUDIT",    "Hash-chained AgentLog (sha256) · tamper-evident", CYAN),
    ]
    y = Inches(1.75)
    for i, (label, body, col) in enumerate(items):
        row_x = Inches(0.55) if i % 2 == 0 else Inches(6.85)
        if i % 2 == 0 and i > 0:
            y += Inches(1.6)
        add_rect(s, row_x, y, Inches(5.95), Inches(1.45),
                 fill=BG_PANEL, line=col)
        add_accent_bar(s, row_x, y, Inches(1.45), col)
        add_text(s, row_x + Inches(0.25), y + Inches(0.15), Inches(5.5), Inches(0.4),
                 label, size=14, bold=True, color=col)
        add_text(s, row_x + Inches(0.25), y + Inches(0.55), Inches(5.5), Inches(0.85),
                 body, size=12, color=WHITE)

    # Footer reference
    add_text(s, Inches(0.55), Inches(6.7), Inches(12.0), Inches(0.4),
             "OWASP Top-10 (2021) — every category mapped to a control. Full table in SECURITY_AUDIT.md",
             size=11, color=GREY, align=PP_ALIGN.CENTER)
    page_number(s, 16)
    set_notes(s,
        "Omar Kapil. 'A security tool that is insecure is worse than no tool. We applied "
        "the same OWASP checks to ourselves.' Direct examiners to SECURITY_AUDIT.md.")

# ── Slide 17 — Lab Environment ────────────────────────────────────────
def slide_17():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "The Lab Environment",
               "Simulated SME · 4 isolated zones · 3 documented attack scenarios", accent=RED)
    zones = [
        ("DMZ", "10.10.10.0/24", RED, [
            ("Juice Shop",   "10.10.10.10", "SQLi · XSS · BOLA · SSRF", "9.5"),
            ("API Gateway",  "10.10.10.20", "Header leak · info disclosure", "6.0"),
            ("CoreDNS",      "10.10.10.30", "Zone xfer · DNS amplification", "5.0"),
        ]),
        ("CORP", "10.10.20.0/24", ORANGE, [
            ("Samba file",   "10.10.20.10", "admin/admin123 · enum",  "8.0"),
            ("GreenMail",    "10.10.20.20", "Plaintext SMTP · POP3",  "7.0"),
            ("Workstation",  "10.10.20.40", "Internal info disclosure","4.0"),
        ]),
        ("DATA", "10.10.30.0/24", MAGENTA, [
            ("PostgreSQL",   "10.10.30.10", "Weak password · no SSL", "9.0"),
            ("Redis cache",  "10.10.30.20", "No auth · protect-mode off", "8.5"),
        ]),
        ("MGMT", "10.10.40.0/24", CYAN, [
            ("traffic_gen",  "10.10.40.10", "Background HTTP/SMB",    "—"),
            ("log_shipper",  "10.10.40.20", "→ Elasticsearch / Wazuh","—"),
        ]),
    ]
    x = Inches(0.55); y = Inches(1.7)
    for i, (zname, cidr, col, hosts) in enumerate(zones):
        row_x = x + Inches((i % 2) * 6.3)
        row_y = y + Inches((i // 2) * 2.55)
        add_rect(s, row_x, row_y, Inches(5.95), Inches(2.4),
                 fill=BG_PANEL, line=col)
        add_text(s, row_x + Inches(0.2), row_y + Inches(0.1), Inches(5.5), Inches(0.4),
                 f"{zname}   {cidr}", size=13, bold=True, color=col, font="Consolas")
        ty = row_y + Inches(0.55)
        for hname, hip, hdesc, cvss in hosts:
            add_text(s, row_x + Inches(0.2), ty, Inches(1.6), Inches(0.4),
                     hname, size=11, bold=True, color=WHITE,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, row_x + Inches(1.85), ty, Inches(1.4), Inches(0.4),
                     hip, size=10, color=col, font="Consolas",
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, row_x + Inches(3.25), ty, Inches(2.3), Inches(0.4),
                     hdesc, size=10, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, row_x + Inches(5.55), ty, Inches(0.4), Inches(0.4),
                     cvss, size=10, bold=True, color=RED if cvss not in ("—",) else GREY,
                     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
            ty += Inches(0.45)
    page_number(s, 17)
    set_notes(s,
        "Shahd Paher. 'This is not a single Juice Shop container — it is a 10-container "
        "miniature SME with realistic background traffic. Every scenario has a CVE, a "
        "CVSS score, and an expected dashboard outcome.'")

# ── Slide 18 — Live Demo ──────────────────────────────────────────────
def slide_18():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    # giant DEMO marker
    big = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(0.7), Inches(11.7), Inches(2.2))
    big.adjustments[0] = 0.08
    big.fill.solid(); big.fill.fore_color.rgb = BG_PANEL
    big.line.color.rgb = CYAN; big.line.width = Pt(2)
    big.shadow.inherit = False
    add_text(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(2.2),
             "LIVE  DEMO",
             size=80, bold=True, color=CYAN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.8), Inches(2.95), Inches(11.7), Inches(0.5),
             "Scan to Action in under 3 minutes",
             size=18, color=GREY, align=PP_ALIGN.CENTER)

    # Cue cards
    cues = [
        ("01", "Login",                "admin@local · dashboard loads"),
        ("02", "Trigger AI scan",      "target = 10.10.10.10 (Juice Shop)"),
        ("03", "Watch Orchestration Feed", "RECON → ATTACK → VALIDATE → RISK"),
        ("04", "Topology turns red",   "node colour by risk score"),
        ("05", "Open Asset Detail",    "AI plain-English advisory"),
        ("06", "Action Center",        "5 prioritised remediation items"),
        ("07", "Reports tab",          "one-click PDF export"),
    ]
    y = Inches(3.7); cw = Inches(1.66); gap = Inches(0.07)
    for i, (n, t, sub) in enumerate(cues):
        cx = Inches(0.55) + (cw + gap) * i
        add_rect(s, cx, y, cw, Inches(2.7), fill=BG_PANEL2, line=CYAN)
        add_accent_bar(s, cx, y, Inches(2.7), CYAN, w=Inches(0.06))
        add_text(s, cx + Inches(0.2), y + Inches(0.15), cw, Inches(0.4),
                 n, size=14, bold=True, color=CYAN)
        add_text(s, cx + Inches(0.2), y + Inches(0.55), cw - Inches(0.3), Inches(0.7),
                 t, size=13, bold=True, color=WHITE)
        add_text(s, cx + Inches(0.2), y + Inches(1.3), cw - Inches(0.3), Inches(1.3),
                 sub, size=10, color=GREY)
    add_text(s, Inches(11.7), Inches(7.05), Inches(1.1), Inches(0.35),
             "18 / 20", size=10, color=CYAN, align=PP_ALIGN.RIGHT)
    set_notes(s,
        "Omar Kapil drives demo from demo/demo_script.md. Yousef on agent log lines, "
        "Shahd on scan internals, Marize/Omnia on UI. Time-cap 10 min. If a step "
        "hangs >30s, switch to pre-recorded fallback video.")

# ── Slide 19 — Team ───────────────────────────────────────────────────
def slide_19():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "How 12 People Built This",
               "Team Lead + 4 sub-team leads + 7 specialists")

    # Team Lead box (top centre)
    lead_x = Inches(4.8); lead_y = Inches(1.7)
    add_rect(s, lead_x, lead_y, Inches(3.7), Inches(0.95),
             fill=BG_PANEL2, line=CYAN, line_w=2)
    add_text(s, lead_x, lead_y + Inches(0.05), Inches(3.7), Inches(0.4),
             "OMAR KAPIL", size=14, bold=True, color=CYAN,
             align=PP_ALIGN.CENTER)
    add_text(s, lead_x, lead_y + Inches(0.45), Inches(3.7), Inches(0.5),
             "Team Lead  ·  DevOps Sub-Leader\n(docker · CI · Caddy · lab)",
             size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # 4 sub-team columns
    teams = [
        ("SUB-TEAM 1",  "Backend & AI Core",  CYAN,
         ("Reem Amin",    "Backend Lead\nFastAPI · auth · DB"),
         [("Yousef Abdel Hady", "Agent pipeline · risk engine · Gemini"),
          ("Mohamed Shaban",    "Celery · Redis · WebSocket · reaper")]),
        ("SUB-TEAM 2",  "Frontend & Visualization", MAGENTA,
         ("Marize Ehap",  "Frontend Lead\nReact · routing · auth UI"),
         [("Omnia Helmy",       "D3 force graph · charts · heatmaps"),
          ("Rahma Ebrahem",     "UI/UX · cyber theme · a11y · toasts")]),
        ("SUB-TEAM 3",  "Security & Scanning", RED,
         ("Shahd Paher",  "Security Lead\nNmap · Nuclei · OpenVAS · lab"),
         [("Mariz Ehap",        "Wazuh · Elasticsearch · Kibana · correlator"),
          ("", "")]),
        ("SUB-TEAM 4",  "DevOps & QA", GREEN,
         ("Omar Kapil",   "DevOps Lead (also Team Lead)\ncompose · CI · TLS"),
         [("Yosef Ali",          "API tests · Postman · coverage"),
          ("Mazin Alla",         "Playwright E2E · UAT · cross-browser")]),
    ]
    base_y = Inches(2.85)
    col_w = Inches(2.95); gap = Inches(0.1)
    start_x = Inches(0.55)
    for i, (label, name, col, lead, members) in enumerate(teams):
        cx = start_x + (col_w + gap) * i
        # team header
        add_rect(s, cx, base_y, col_w, Inches(0.55),
                 fill=col, line=col)
        add_text(s, cx, base_y, col_w, Inches(0.55),
                 label, size=11, bold=True, color=BG_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # team name
        add_rect(s, cx, base_y + Inches(0.55), col_w, Inches(0.45),
                 fill=BG_PANEL2, line=col)
        add_text(s, cx, base_y + Inches(0.55), col_w, Inches(0.45),
                 name, size=11, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # lead
        ly = base_y + Inches(1.05)
        add_rect(s, cx, ly, col_w, Inches(0.95),
                 fill=BG_PANEL, line=col, line_w=1.5)
        add_text(s, cx + Inches(0.1), ly + Inches(0.05), col_w, Inches(0.35),
                 lead[0], size=12, bold=True, color=WHITE)
        add_text(s, cx + Inches(0.1), ly + Inches(0.4), col_w - Inches(0.2), Inches(0.55),
                 lead[1], size=9, color=col)
        # members
        my = ly + Inches(1.05)
        for mname, mtask in members:
            if not mname:
                continue
            add_rect(s, cx, my, col_w, Inches(0.85),
                     fill=BG_PANEL, line=LINE)
            add_text(s, cx + Inches(0.1), my + Inches(0.05), col_w, Inches(0.32),
                     mname, size=11, bold=True, color=WHITE)
            add_text(s, cx + Inches(0.1), my + Inches(0.35), col_w - Inches(0.2), Inches(0.5),
                     mtask, size=9, color=GREY)
            my += Inches(0.95)

    # Cross-team contributor (Omar Tarek)
    omar_t_y = Inches(6.5)
    add_rect(s, Inches(0.55), omar_t_y, Inches(12.2), Inches(0.55),
             fill=BG_PANEL2, line=ORANGE)
    add_text(s, Inches(0.75), omar_t_y, Inches(12.0), Inches(0.55),
             "OMAR TAREK  —  Documentation & Presentation Lead   "
             "(Mermaid diagrams · demo script · slides · video · Swagger)",
             size=12, bold=True, color=ORANGE,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    page_number(s, 19)
    set_notes(s,
        "Omar Kapil. Walk top-to-bottom. Team Lead at top; 4 sub-team leads as the "
        "second tier; specialists under each lead. ONE LINE per person. ~60 seconds "
        "total for the slide.")

# ── Slide 20 — Conclusion / Q&A ───────────────────────────────────────
def slide_20():
    s = prs.slides.add_slide(BLANK)
    add_header(s, "Conclusion · Future Work · Q&A")
    # Done
    add_rect(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(5.0),
             fill=BG_PANEL, line=GREEN)
    add_text(s, Inches(0.85), Inches(1.85), Inches(5.5), Inches(0.4),
             "WHAT WE ACHIEVED", size=12, bold=True, color=GREEN)
    add_bullets(s, Inches(0.85), Inches(2.3), Inches(5.5), Inches(4.3), [
        "40+ documented REST endpoints",
        "4-stage AI agent pipeline (production-grade)",
        "Deterministic 0–100 risk engine",
        "25+ dashboard panels · live WebSocket UI",
        "4-zone simulated SME lab · 3 attack scenarios",
        "Wazuh + Elasticsearch SIEM integration",
        "≥ 15 backend tests + Playwright E2E suite",
        "One-command Docker deploy under 4 minutes",
    ], size=13, bullet_color=GREEN)

    # Roadmap
    add_rect(s, Inches(6.85), Inches(1.7), Inches(5.95), Inches(5.0),
             fill=BG_PANEL, line=CYAN)
    add_text(s, Inches(7.05), Inches(1.85), Inches(5.5), Inches(0.4),
             "WHAT'S NEXT", size=12, bold=True, color=CYAN)
    add_bullets(s, Inches(7.05), Inches(2.3), Inches(5.5), Inches(4.3), [
        "Cosign image signing (supply chain)",
        "Cloud asset discovery (AWS / GCP)",
        "n8n SOAR playbooks (auto-remediation)",
        "Celery worker autoscaling",
        "Multi-tenant deployment for MSPs",
    ], size=13, bullet_color=CYAN)

    # Footer banner
    add_rect(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.45),
             fill=BG_PANEL2, line=CYAN)
    add_text(s, Inches(0.75), Inches(6.85), Inches(12.0), Inches(0.45),
             "Thank you  —  questions are welcome  ·  HOW_TO_RUN.md reproduces the demo in <5 min",
             size=12, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    page_number(s, 20)
    set_notes(s,
        "Omar Kapil keeps time and re-routes Q&A: Reem/Shahd/Marize for technical "
        "questions, Yousef for AI questions. Q&A prepared answers are in "
        "FINAL_PRESENTATION.md.")

# ── Build ─────────────────────────────────────────────────────────────
for fn in [slide_01, slide_02, slide_03, slide_04, slide_05, slide_06, slide_07,
           slide_08, slide_09, slide_10, slide_11, slide_12, slide_13, slide_14,
           slide_15, slide_16, slide_17, slide_18, slide_19, slide_20]:
    fn()

OUT = "FINAL_PRESENTATION.pptx"
prs.save(OUT)
print(f"Saved {OUT}  ·  {len(prs.slides)} slides")
